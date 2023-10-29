"""
Add automatically alt text to each picture in a powerpoint file
"""

import os
import sys
import argparse
import requests
import base64
from typing import List
from pptx import Presentation
from pptx.shapes.base import BaseShape

def check_server_is_running(url: str) -> bool:
    try:
        response = requests.get(url)
        if response.status_code == 200:
            return True
    except requests.ConnectionError:
        return False
    return False

# https://github.com/scanny/python-pptx/pull/512
def shape_get_alt_text(shape: BaseShape) -> str:
    """ Alt-text defined in shape's `descr` attribute, or '' if not present. """
    return shape._element._nvXxPr.cNvPr.attrib.get("descr", "")

def shape_set_alt_text(shape: BaseShape, alt_text: str):
    """ Set alt-text in shape """
    shape._element._nvXxPr.cNvPr.attrib["descr"] = alt_text

def process_images_from_pptx(file_path: str, set_image_description: bool, server_url: str, prompt: str, DEBUG: bool = False) -> None:
    """ 
    Loop through images in the slides of a Powerpint file and set image description based 
    on image description from Llava
    """
    file_name = os.path.basename(file_path)

    # get name, extension, folder
    name = file_name.split(".")[0]
    extension = file_name.split(".")[1]
    dirname = os.path.dirname(file_path)

    # Initialize presentation object
    print(f"Reading {file_path}")
    prs = Presentation(file_path)
    
    # Write alt-text to file
    out_file_name = os.path.join(dirname, f"{name}.txt")
    fout = open(out_file_name, "w")
    # write header
    fout.write(f"Powerpoint\tSlide\tPicture\tAlt_Text{os.linesep}")

    # Loop through slides
    slide_cnt = 1
    for slide in prs.slides:

        for shape in slide.shapes:

            # Check if the shape has a picture
            if shape.shape_type == 13:  # Shape type 13 corresponds to a picture
                
                if set_image_description:

                    image_file_name = f"slide_{slide_cnt}_pict_{shape.name}.jpg"
                    image_file_path = os.path.join("tmp", image_file_name)
                    print(f"Processing image: {image_file_path}")
                    
                    # get image
                    image_stream = shape.image.blob

                    # save image
                    with open(image_file_path, "wb") as f:
                        f.write(image_stream)

                    # read image
                    with open(image_file_path, 'rb') as img_file:
                        img_byte_arr = img_file.read()

                    # encode in base64
                    img_base64 = base64.b64encode(img_byte_arr).decode('utf-8')

                    # Use LLaVa to get image descriptions
                    header = {"Content-Type": "application/json"}
                    data = {
                        "image_data": [{"data": img_base64, "id": 1}],
                        "prompt": f"USER:[img-1] {prompt}\nASSISTANT:",
                        "n_predict": 123,
                        "temp": 0.1
                    }
                    response = requests.post(server_url, headers=header, json=data)
                    response_data = response.json()

                    if DEBUG:
                        print(response_data)
                        print()

                    alt_text = response_data.get('content', '').strip()
                    if DEBUG:
                        print(f"Len: {len(alt_text)}, Content: {alt_text}")

                    if len(response.text) > 0:
                        image_description = alt_text
                        shape_set_alt_text(shape, image_description)
                    else:
                        print("No content.")

                # report alt text
                stored_alt_text = shape_get_alt_text(shape)
                feedback = f"Slide: {slide_cnt}, Picture: '{shape.name}', alt_text: '{stored_alt_text}'"
                print(feedback)
                fout.write(f"{name}.{extension}\t{slide_cnt}\t{shape.name}\t{stored_alt_text}" + os.linesep)

        slide_cnt += 1

    fout.close()

    print(f"Powerpoint file contains {slide_cnt} slides.")

    if set_image_description:
        # Save file
        outfile = os.path.join(dirname, f"{name}_alt_text.{extension}")
        print(f"Saving to {outfile}")
        prs.save(outfile)

def main(argv: List[str]) -> int:
    """
    Add alt-text automatically to images in PowerPoint.

    Parameters:
        argv (List[str]): Command line arguments

    Returns:
        int: Exit code (0 if successful, else non-zero)
    """    
    parser = argparse.ArgumentParser(description='Add alt-text automatically to images in Powerpoint')
    parser.add_argument("file", type=str, help="Powerpoint file")
    parser.add_argument("--add", action='store_true', default=False, help="Flag to add alt-text to images")
    parser.add_argument("--prompt", type=str, default="Describe image clearly and detailed. \
                        Check if graph and if so report summary of what the graph depicts. \
                        Make sure to ouput only up to 125 characters.", help="LLaVA prompt")
    parser.add_argument("--debug", action='store_true', default=False, help="Debug")

    args = parser.parse_args()

    server_url = "http://localhost:8007"
    if check_server_is_running(server_url):
        server_url = "http://localhost:8007/completion"
        # Read PowerPoint file and list images
        powerpoint_file_name = args.file
        
        if os.path.isfile(powerpoint_file_name):
            process_images_from_pptx(powerpoint_file_name, args.add, server_url, args.prompt, args.debug)
        else:
            print(f"Error: File {powerpoint_file_name} not found.")
            return 1
        
        return 0
    else:
        print(f"Unable to access server at '{server_url}'")
        return 1

if __name__ == "__main__":
    exit_code = main(sys.argv[1:])
    sys.exit(exit_code)
