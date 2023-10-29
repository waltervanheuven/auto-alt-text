from pptx import Presentation
from pptx.shapes.base import BaseShape
import os

# https://github.com/scanny/python-pptx/pull/512
def shape_get_alt_text(shape: BaseShape) -> str:
    """ Alt-text defined in shape's `descr` attribute, or '' if not present. """
    return shape._element._nvXxPr.cNvPr.attrib.get("descr", "")

def shape_set_alt_text(shape: BaseShape, alt_text: str):
    """ Set alt-text in shape """
    shape._element._nvXxPr.cNvPr.attrib["descr"] = alt_text

def process_images_from_pptx(file_path: str, set_image_description: bool) -> None:
    """ 
    Loop through images in the slides of a Powerpint file and set image description based 
    on image description from Llava
    """
    file_name = os.path.basename(file_path)

    # get name, extension, folder
    name = file_name.split(".")[0]
    extension = file_name.split(".")[1]
    dirname = os.path.dirname(file_path)

    # Initialize presentation object
    print(f"Reading {file_path}")
    prs = Presentation(file_path)
        
    # Loop through slides
    slide_cnt = 1
    for slide in prs.slides:

        for shape in slide.shapes:

            # Check if the shape has a picture
            if shape.shape_type == 13:  # Shape type 13 corresponds to a picture
                
                if set_image_description:
                    # get image
                    image_stream = shape.image.blob

                    # save image                
                    image_file_name = f"slide_{slide_cnt}_pict_{shape.name}.jpg"
                    image_file_path = os.path.join("tmp", image_file_name)
                    with open(image_file_path, "wb") as f:
                        f.write(image_stream)

                    # Use LLaVa to get image descriptions
                    image_description = ""
                    shape_set_alt_text(shape, image_description)

                # report alt text
                print(f"Slide: {slide_cnt}, Picture: '{shape.name}', alt_text: '{shape_get_alt_text(shape)}'")

        slide_cnt += 1

    print(f"Powerpoint file contains {slide_cnt} slides.")

    if set_image_description:
        # Save file
        outfile = os.path.join(dirname, f"{name}_alt_text.{extension}")
        print(f"Saving to {outfile}")
        prs.save(outfile)

# Read PowerPoint file and list images
powerpoint_file_name = "tmp/test.pptx"
process_images_from_pptx(powerpoint_file_name, False)
