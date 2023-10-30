"""
Add automatically alt text to each picture in a powerpoint file using OpenCLIP
"""

import os
import sys
import argparse
from typing import List
from pptx import Presentation
from pptx.shapes.base import BaseShape
import open_clip
import torch
from PIL import Image

def num2str(the_max: int, n:int) -> str:
    if the_max > 99:
        if n < 100:
            if n < 10:
                return f"00{str(n)}"
            else:
                return f"0{str(n)}"
        else:
            return f"{str(n)}"
    else:
        if n < 10:
            return f"0{str(n)}"
        else:
            return f"{str(n)}"

# https://github.com/scanny/python-pptx/pull/512
def shape_get_alt_text(shape: BaseShape) -> str:
    """ Alt-text defined in shape's `descr` attribute, or '' if not present. """
    return shape._element._nvXxPr.cNvPr.attrib.get("descr", "")

def shape_set_alt_text(shape: BaseShape, alt_text: str):
    """ Set alt-text in shape """
    shape._element._nvXxPr.cNvPr.attrib["descr"] = alt_text

def process_images_from_pptx(file_path: str, set_image_description: bool, model_name: str, pretrained: str, DEBUG: bool = False) -> bool:
    """ 
    Loop through images in the slides of a Powerpint file and set image description based 
    on image description from OpenClip
    """
    err: bool = False
    file_name:str = os.path.basename(file_path)

    # get name, extension, folder
    name:str = file_name.split(".")[0]
    extension:str = file_name.split(".")[1]
    dirname:str = os.path.dirname(file_path)

    # Initialize presentation object
    print(f"Reading {file_path}")
    prs = Presentation(file_path)
    
    # Write alt-text to file
    out_file_name:str = os.path.join(dirname, f"{name}.txt")
    fout = open(out_file_name, "w")
    # write header
    fout.write(f"Powerpoint\tSlide\tPicture\tAlt_Text{os.linesep}")

    nr_slides = len(prs.slides)

    # OpenCLIP model
    if set_image_description:
        model, _, transform = open_clip.create_model_and_transforms(
            model_name=model_name,
            pretrained=pretrained
        )

    # Loop through slides
    slide_cnt:int = 1
    image_cnt:int = 1
    for slide in prs.slides:
        # loop through shapes
        slide_image_cnt = 1
        for shape in slide.shapes:
            # Check if the shape has a picture
            if shape.shape_type == 13:  # Shape type 13 corresponds to a picture
                
                if set_image_description:
                    err = set_alt_text(shape, slide_cnt, nr_slides, slide_image_cnt, model, transform, DEBUG)
                
                # report alt text
                if not err:
                    stored_alt_text:str = shape_get_alt_text(shape)
                    feedback = f"Slide: {slide_cnt}, Picture: '{shape.name}', alt_text: '{stored_alt_text}'"
                    print(feedback)
                    fout.write(f"{name}.{extension}\t{slide_cnt}\t{shape.name}\t{stored_alt_text}" + os.linesep)

                    slide_image_cnt += 1
                    image_cnt += 1

        slide_cnt += 1

    fout.close()

    slide_cnt -= 1
    image_cnt -= 1
    print(f"Powerpoint file contains {slide_cnt} slides and {image_cnt} images.")

    if set_image_description:
        # Save file
        outfile:str = os.path.join(dirname, f"{name}_alt_text.{extension}")
        print(f"Saving to {outfile}")
        prs.save(outfile)

    return err

def set_alt_text(shape: BaseShape, slide_cnt: int, max_slides: int, image_cnt: int, model, transform, DEBUG: bool) -> bool:
    err: bool = False
    
    # get image
    try:
        image_stream = shape.image.blob
        extension:str = shape.image.ext
    except Exception as e:
        print(f"Exception {str(e)}")
        err = True

    if not err:
        image_file_name = f"s{num2str(max_slides, slide_cnt)}p{num2str(99, image_cnt)}_{shape.name}"
        image_file_path = os.path.join("tmp", image_file_name)
        image_file_path = os.path.join("tmp", f"{image_file_name}.{extension}")
        print(f"Processing image: '{image_file_path}'...")

        # save image
        with open(image_file_path, "wb") as f:
            f.write(image_stream)

        # read image
        im = Image.open(image_file_path).convert("RGB")
        im = transform(im).unsqueeze(0)

        with torch.no_grad(), torch.cuda.amp.autocast():
            generated = model.generate(im)

        # get picture description and remove trailing spaces
        alt_text = open_clip.decode(generated[0]).split("<end_of_text>")[0].replace("<start_of_text>", "")

        if DEBUG:
            print(f"Len: {len(alt_text)}, Content: {alt_text}")

        if len(alt_text) > 0:
            image_description = alt_text
            shape_set_alt_text(shape, image_description)
        else:
            print("No content.")

    return err

def main(argv: List[str]) -> int:
    err: bool = False

    parser = argparse.ArgumentParser(description='Add alt-text automatically to images in Powerpoint')
    parser.add_argument("file", type=str, help="Powerpoint file")
    parser.add_argument("--add", action='store_true', default=False, help="flag to add alt-text to images")
    parser.add_argument("--model", type=str, default="coca_ViT-L-14", help="model name")
    parser.add_argument("--pretrained", type=str, default="mscoco_finetuned_laion2B-s13B-b90k", help="pretrained model")
    parser.add_argument("--debug", action='store_true', default=False, help="debug")

    args = parser.parse_args()

    # Read PowerPoint file and list images
    powerpoint_file_name = args.file
    if not os.path.isfile(powerpoint_file_name):
        print(f"Error: File {powerpoint_file_name} not found.")
        err = True
    
    if not err:
        err = process_images_from_pptx(powerpoint_file_name, args.add, args.model, args.pretrained, args.debug)

    if err:
        return 1
    else:
        return 0

if __name__ == "__main__":
    exit_code = main(sys.argv[1:])
    sys.exit(exit_code)
