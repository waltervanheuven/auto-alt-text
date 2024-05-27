""" powerpoint.py """

from typing import Tuple
import os
import sys
import io
import platform
import pathlib
import shutil
import subprocess
from PIL import Image
from pptx.util import Cm
from pptx.oxml.ns import _nsmap
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape
from pptx import Presentation
from .utils import num2str, str2bool, bool2str
from .models import kosmos2, openclip, qwen_vl, cog_vlm, phi3_vision
from .models import use_ollama, use_openai, use_mlx_vlm

# see https://github.com/scanny/python-pptx/pull/512
def get_alt_text(shape: BaseShape) -> str:
    """ Alt text is defined in shape's `descr` attribute, return this or '' if not present. """
    return shape._element._nvXxPr.cNvPr.attrib.get("descr", "")

def set_alt_text(shape: BaseShape, alt_text: str) -> None:
    """ Set alt text of shape """
    try:
        shape._element._nvXxPr.cNvPr.attrib["descr"] = alt_text
    except Exception as e:
        print(f"--> Unable to set alt_text: {shape.shape_type}, {shape.name}\n{str(e)}\nAlt_text: {alt_text}")

# see https://stackoverflow.com/questions/63802783/check-if-image-is-decorative-in-powerpoint-using-python-pptx
def is_decorative(shape) -> bool:
    """ check if image is decorative """
    # <adec:decorative xmlns:adec="http://schemas.microsoft.com/office/drawing/2017/decorative" val="1"/>
    _nsmap["adec"] = "http://schemas.microsoft.com/office/drawing/2017/decorative"
    cNvPr = shape._element._nvXxPr.cNvPr
    adec_decoratives = cNvPr.xpath(".//adec:decorative[@val='1']")
    return bool(adec_decoratives)

def process_shape(
        shape: BaseShape,
        pptx: dict,
        settings: dict,
        verbose: bool = False,
        debug: bool = False
    ) -> bool:
    """
    Recursive function to process shapes and shapes in groups on each slide
    """
    err: bool = False
    report:bool = settings["report"]

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        # keep a list of images as part of the group
        if pptx["group_shape_list"] is None:
            pptx["group_shape_list"] = [shape]
        else:
            group_shape_list = pptx["group_shape_list"]
            group_shape_list.append(shape)
            pptx["group_shape_list"] = group_shape_list

        if pptx["image_list"] is None:
            pptx["image_list"] = []
        if pptx["text_list"] is None:
            pptx["text_list"] = []
        if pptx["object_list"] is None:
            pptx["object_list"] = []

        # process shapes
        for embedded_shape in shape.shapes:
            err = process_shape(embedded_shape, pptx, settings, debug)
            if err:
                break

        if not err:
            # check if group is not part of other group
            group_shape_list:list[BaseShape] = pptx["group_shape_list"]
            part_of_group:str = "No"
            if len(group_shape_list) > 1:
                part_of_group = "Yes"
            elif len(group_shape_list) == 1:
                part_of_group = "No_TopLevel"

            # current group shape
            group_shape:BaseShape = get_current_group_shape(pptx)

            # image list
            image_list:list = pptx["image_list"]

            # group contains at least one image
            # if image_list is not None and len(image_list) > 0:
            #     new_img = combine_images_in_group(image_list, group_shape)
            #     img_folder = pptx["img_folder"]
            #     filename = os.path.join(img_folder, f'slide_{pptx["slide_cnt"]}_group.png')
            #     new_img.save(filename)

            alt_text:str = ""
            if not report:
                # combine text box content associated with group
                text_list:list = pptx["text_list"]
                for n, txt in enumerate(text_list):
                    # remove newlines
                    txt = replace_newline_with_space(txt)
                    if n == 0:
                        alt_text = txt
                    else:
                        alt_text = f"{alt_text} {txt}"

                if len(alt_text) > 0:
                    alt_text = f"{alt_text}. "

                # combine alt text to generate the alt text for the group

                if len(image_list) > 1:
                    alt_text = f"{alt_text}There are {len(image_list)} images:"
                for _, _, _, txt in image_list:
                    # remove newlines
                    txt = replace_newline_with_space(txt)
                    if len(alt_text) == 0:
                        alt_text = txt
                    else:
                        alt_text = f"{alt_text} {txt}"

                # set alt text of group shape
                set_alt_text(group_shape, alt_text)
            else:
                alt_text = get_alt_text(group_shape)

                # remove returns
                alt_text = replace_newline_with_space(alt_text)

            # get vars
            model_str:str = settings["model"]
            report:bool = settings["report"]
            image_file_path:str = ""
            pptx_name:str = pptx["pptx_name"]
            pptx_extension:str = pptx["pptx_extension"]
            slide_cnt:int = pptx["slide_cnt"]

            # get info from groupshape
            decorative:bool = is_decorative(group_shape)
            stored_alt_text:str = get_alt_text(group_shape)

            if verbose:
                if decorative:
                    print(f"Slide: {slide_cnt + 1}, Group: {group_shape.name}, alt_text: '{stored_alt_text}', decorative")
                else:
                    print(f"Slide: {slide_cnt + 1}, Group: {group_shape.name}, alt_text: '{stored_alt_text}'")

            fout = pptx["fout"]
            fout.write(f"{model_str}\t{pptx_name}{pptx_extension}\t{slide_cnt + 1}\t{group_shape.name}\tGroup\t{part_of_group}\t{stored_alt_text}\t{len(stored_alt_text)}\t{bool2str(decorative)}\t{image_file_path}\n")

            # remove last one
            group_shape_list = pptx["group_shape_list"]
            pptx["group_shape_list"] = group_shape_list[:-1]

    elif shape.shape_type in [MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE]:
        # picture
        image_file_path:str = ""
        decorative:bool = is_decorative(shape)
        group_shape:BaseShape = get_current_group_shape(pptx)

        # only generate alt text when generate options is True and decorative is False
        if not decorative:
            err, image_file_path = process_shape_and_generate_alt_text(shape, pptx, settings, debug)

        if not err:
            part_of_group = "No"
            if group_shape is not None:
                part_of_group = "Yes"

            # report alt text
            if not err:
                slide_cnt:int = pptx["slide_cnt"]
                slide_image_cnt:int = pptx["slide_image_cnt"]

                stored_alt_text = get_alt_text(shape)
                if verbose:
                    if decorative:
                        print(f"Slide: {slide_cnt + 1}, Pict: {slide_image_cnt + 1}, alt_text: '{stored_alt_text}', decorative")
                    else:
                        print(f"Slide: {slide_cnt + 1}, Pict: {slide_image_cnt + 1}, alt_text: '{stored_alt_text}'")

                model_str:str = settings["model"]
                pptx_name:str = pptx["pptx_name"]
                pptx_extension:str = pptx["pptx_extension"]
                fout = pptx["fout"]
                fout.write(f"{model_str}\t{pptx_name}{pptx_extension}\t{slide_cnt + 1}\t{shape.name}\tPicture\t{part_of_group}\t{stored_alt_text}\t{len(stored_alt_text)}\t{bool2str(decorative)}\t{image_file_path}\n")

                pptx["slide_image_cnt"] = slide_image_cnt + 1

    elif shape.shape_type in [MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.FREEFORM, \
                              MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.IGX_GRAPHIC, MSO_SHAPE_TYPE.CANVAS, \
                              MSO_SHAPE_TYPE.MEDIA, MSO_SHAPE_TYPE.WEB_VIDEO, MSO_SHAPE_TYPE.DIAGRAM, \
                              MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT, MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT, \
                              MSO_SHAPE_TYPE.WEB_VIDEO, MSO_SHAPE_TYPE.LINKED_OLE_OBJECT, \
                              MSO_SHAPE_TYPE.CALLOUT]:

        process_object(shape, pptx, settings, debug)

        object_list = pptx["object_list"]
        if object_list is None:
            object_list = []
            object_list.append(shape)
            pptx["object_list"] = object_list

    elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
        # TEXT_BOX is part of a group
        # For the Alt Text it would be useful to add this text to the Alt Text
        group_shape:BaseShape = get_current_group_shape(pptx)
        if group_shape is not None:
            text:str = shape.text_frame.text
            text_list = pptx["text_list"]
            if text_list is not None:
                text_list.append(text)
                pptx["text_list"] = text_list
    elif debug:
        print(f"=> OBJECT: {shape.name}, type: {shape.shape_type}")

    return err

def replace_newline_with_space(txt: str) -> str:
    """ replace newline with space and replace tab with space """
    s = " ".join(txt.splitlines())
    return s.replace("\t", " ")

def get_current_group_shape(pptx:dict) -> BaseShape:
    """ return group shape """
    group_shape_list:list[BaseShape] = pptx["group_shape_list"]
    if group_shape_list is not None and len(group_shape_list) > 0:
        return group_shape_list[-1]

    return None

def shape_type2str(shape_type) -> str:
    """ return name of shape as string """
    s:str = ""

    if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        s =  "Auto shape"
    elif shape_type == MSO_SHAPE_TYPE.LINE:
        s =  "Line"
    elif shape_type == MSO_SHAPE_TYPE.IGX_GRAPHIC:
        s =  "IgxGraphic"
    elif shape_type == MSO_SHAPE_TYPE.CHART:
        s = "Chart"
    elif shape_type == MSO_SHAPE_TYPE.FREEFORM:
        s = "FreeForm"
    elif shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
        s = "TextBox"
    elif shape_type == MSO_SHAPE_TYPE.CANVAS:
        s = "Canvas"
    elif shape_type == MSO_SHAPE_TYPE.MEDIA:
        s = "Media"
    elif shape_type == MSO_SHAPE_TYPE.WEB_VIDEO:
        s = "WebVideo"
    elif shape_type == MSO_SHAPE_TYPE.DIAGRAM:
        s = "Diagram"
    elif shape_type == MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT:
        s = "Control object"
    elif shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
        s = "Embedded object"
    elif shape_type == MSO_SHAPE_TYPE.WEB_VIDEO:
        s = "Web video"

    return s

def process_object(
        shape: BaseShape,
        pptx: dict,
        settings: dict,
        verbose: bool = False,
        debug:bool = False
    ) -> None:
    """ process """
    # only include if it is not part of a group
    # Powerpoint only reports an accessibility error for a missing group shape alt text
    image_file_path:str = ""
    decorative:bool = is_decorative(shape)
    report:bool = settings["report"]

    # include all text inside shape?
    include_all_paragraphs = True

    group_shape:BaseShape = get_current_group_shape(pptx)
    part_of_group:str = "No"
    if group_shape is not None:
        part_of_group = "Yes"

    alt_text:str = ""
    if not report:
        # Quick fix for alt text of shapes
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if len(shape.name) > 0:
                alt_text = f"{cleanup_name_object(shape.name)}"
            else:
                alt_text = ""
        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            if len(shape.name) > 0:
                if cleanup_name_object(shape.name.lower()) != "chart":
                    # avoid duplication chart
                    alt_text = f"{cleanup_name_object(shape.name)}"
                else:
                    alt_text = f"{cleanup_name_object(shape.name)}"
            else:
                alt_text = ""

            # add title of chart to alt_text
            the_chart = shape.chart
            if the_chart.has_title and len(the_chart.chart_title.text_frame.text.strip()) > 0:
                alt_text = f"{alt_text} entitled '{the_chart.chart_title.text_frame.text}'"

        elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
            if len(shape.name) > 0:
                alt_text = f"{cleanup_name_object(shape.name)}"
            else:
                alt_text = ""
        elif shape.shape_type == MSO_SHAPE_TYPE.CANVAS:
            if len(shape.name) > 0:
                alt_text = f"{cleanup_name_object(shape.name)}"
            else:
                alt_text = ""
        elif shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
            if len(shape.name) > 0:
                alt_text = f"{cleanup_name_object(shape.name)}"
            else:
                alt_text = ""
        elif shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
            if len(shape.name) > 0:
                alt_text = f"{cleanup_name_object(shape.name)}"
            else:
                alt_text = ""
        elif shape.shape_type == MSO_SHAPE_TYPE.WEB_VIDEO:
            if len(shape.name) > 0:
                alt_text = f"{cleanup_name_object(shape.name)}"
            else:
                alt_text = ""
        elif shape.shape_type == MSO_SHAPE_TYPE.DIAGRAM:
            if len(shape.name) > 0:
                alt_text = f"{shape.name}"
            else:
                alt_text = ""
        elif shape.shape_type == MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT:
            if len(shape.name) > 0:
                alt_text = f"{shape.name}"
            else:
                alt_text = ""
        elif shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
            if len(shape.name) > 0:
                alt_text = f"{shape.name}"
            else:
                alt_text = ""
        else:
            if len(shape.name) > 0:
                alt_text = f"{shape.name.lower()}"
            else:
                alt_text = ""

        # indicate the text inside the shape
        if shape.has_text_frame:
            if len(shape.text_frame.paragraphs) == 1 and shape.text_frame.paragraphs[0].text != "":
                alt_text = f"{alt_text} with inside the text: {remove_newlines(shape.text_frame.paragraphs[0].text).strip()}"
            else:
                if not include_all_paragraphs:
                    alt_text = f"{alt_text} with text inside."
                else:
                    first = True
                    for p in shape.text_frame.paragraphs:
                        if p.text != "":
                            if first:
                                alt_text = f"{alt_text} with inside the text:"
                                first = False

                            alt_text = f"{alt_text} {remove_newlines(p.text).strip()}"

        # make sure alt text ends with a final stop
        if not alt_text.endswith("."):
            alt_text = f"{alt_text}."

        set_alt_text(shape, alt_text)

        # if part of group store alt_text
        if group_shape is not None:
            text_list = pptx["text_list"]
            if text_list is not None:
                text_list.append(alt_text)
                pptx["text_list"] = text_list

    # report alt text
    slide_cnt:int = pptx["slide_cnt"]
    slide_image_cnt:int = pptx["slide_image_cnt"]

    stored_alt_text = get_alt_text(shape)
    if verbose:
        if decorative:
            print(f"Slide: {slide_cnt + 1}, {shape_type2str(shape.shape_type)}: {slide_image_cnt + 1}, alt_text: '{stored_alt_text}', decorative")
        else:
            print(f"Slide: {slide_cnt + 1}, {shape_type2str(shape.shape_type)}: {slide_image_cnt + 1}, alt_text: '{stored_alt_text}'")

    model_str:str = settings["model"]
    pptx_name:str = pptx["pptx_name"]
    pptx_extension:str = pptx["pptx_extension"]
    fout = pptx["fout"]
    fout.write(f"{model_str}\t{pptx_name}{pptx_extension}\t{slide_cnt + 1}\t{shape.name}\t{shape_type2str(shape.shape_type)}\t{part_of_group}\t{stored_alt_text}\t{len(stored_alt_text)}\t{bool2str(decorative)}\t{image_file_path}\n")

def remove_newlines(txt:str) -> str:
    """ remove newlines and replace tabs with spaces """
    s = "".join(txt.splitlines())
    return s.replace("\t", " ")

def cleanup_name_object(txt:str) -> str:
    """
    check if alt shape name contains a number at the end 
    e.g. "oval 1", "oval 2" and remove the number
    """
    s:str = ""

    elements:list[str] = txt.lower().split()
    if len(elements) == 1:
        s = elements[0]
    else:
        last_word = elements[-1]
        try:
            _ = int(last_word)
            s = ' '.join(elements[:-1])
        except ValueError:
            s = txt

    return s

def combine_images_in_group(
        images, 
        group_shape,
        verbose: bool = False
    ) -> Image:
    """ 
    Create new image based on shape

    TODO: Not yet working properly, image size is not correct
    """

    # EMU per Pixel estimate: not correct
    emu_per_pixel:int = int(914400 / 300)

    # Determine the size of the new image based on the group shape size
    new_img_width = int(group_shape.width / emu_per_pixel)
    new_img_height = int(group_shape.height / emu_per_pixel)
    new_img = Image.new('RGB', (new_img_width, new_img_height))

    # Paste each image into the new image at its relative position
    for image, left, top, _ in images:
        if verbose:
            print(f"img: {image.width}, {image.height}, {left}, {top}")
        new_img.paste(image, (int(left / emu_per_pixel), int(top / emu_per_pixel)))

    return new_img

def process_shape_and_generate_alt_text(
        shape: BaseShape,
        pptx: dict,
        settings: dict,
        verbose: bool = False,
        debug:bool=False
    ) -> Tuple[bool, str]:
    """ 
    Save image associated with shape and generate alt text
    """
    err:bool = False
    image_file_path:str = ""

    # get image
    image_stream = None
    extension:str = ""
    if hasattr(shape, "image"):
        # get image, works with only with PNG, JPEG?
        image_stream = shape.image.blob
        extension = shape.image.ext
    else:
        # get image for other formats, e.g. TIFF
        # <Element {http://schemas.openxmlformats.org/presentationml/2006/main}pic at 0x15f2d6b20>
        try:
            slide_part = shape.part
            rId = shape._element.blip_rId
            image_part = slide_part.related_part(rId)
            image_stream = image_part.blob
            extension = image_part.partname.ext
        except AttributeError:
            slide_cnt:int = pptx["slide_cnt"] + 1
            print(f"Error, slide: {slide_cnt}, pict: '{shape.name}', unable to access image", file=sys.stderr)
            #err = True

    if not err and image_stream is not None:
        report:bool = settings["report"]
        base_left:int = pptx["base_left"]
        base_top:int = pptx["base_top"]

        # determine file name
        pptx_nslides: int = pptx["pptx_nslides"]
        slide_image_cnt:int = pptx["slide_image_cnt"]
        slide_cnt:int = pptx["slide_cnt"]
        pptx_nslides:int = pptx["pptx_nslides"]
        img_folder:str = pptx["img_folder"]

        alt_text:str = ""
        if not report:
            image_file_name:str = f"s{num2str(pptx_nslides, slide_cnt + 1)}p{num2str(99, slide_image_cnt + 1)}"
            image_file_path = os.path.join(img_folder, f"{image_file_name}.{extension}")
            if verbose:
                print(f"Saving image from pptx: '{image_file_path}'")

            # save image
            with open(image_file_path, "wb") as f:
                f.write(image_stream)

            alt_text, err = generate_description(image_file_path, extension, settings, verbose=verbose)
        else:
            alt_text = get_alt_text(shape)

            # remove returns
            alt_text = replace_newline_with_space(alt_text)

        if not err:
            # Keep image in case the image is part of a group
            if pptx["image_list"] is not None:
                # Calculate the position of the image relative to the group
                image_group_part = Image.open(io.BytesIO(image_stream))
                left = base_left + shape.left
                top = base_top + shape.top
                image_list = pptx["image_list"]
                image_list.append((image_group_part, left, top, alt_text))
                pptx["image_list"] = image_list

            if debug:
                print(f"Len: {len(alt_text)}, Content: {alt_text}")

            if len(alt_text) > 0:
                set_alt_text(shape, alt_text)

    return err, image_file_path

def process_shapes_from_file(
        shape: BaseShape,
        group_shape_list: list[BaseShape],
        csv_rows, slide_cnt: int,
        slide_object_cnt: int,
        object_cnt: int,
        verbose: bool = False,
        debug: bool = False
    ) -> Tuple[list[BaseShape], int, int]:
    """ recursive function to process shapes and shapes within groups """
    # Check if the shape has a picture
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        if group_shape_list is None:
            group_shape_list = [shape]
        else:
            group_shape_list.append(shape)

        for embedded_shape in shape.shapes:
            group_shape_list, object_cnt, slide_object_cnt = process_shapes_from_file(embedded_shape, group_shape_list, csv_rows, slide_cnt, slide_object_cnt, object_cnt, verbose, debug)

        # current group shape (last one)
        group_shape = group_shape_list[-1]

        # get decorative
        decorative_pptx:bool = is_decorative(group_shape)
        decorative:bool = str2bool(csv_rows[object_cnt][8])

        # change decorative status
        if decorative_pptx != decorative:
            # set decorative status of image
            print(f"Side: {slide_cnt}, {group_shape.name}, can't set the docorative status to: {bool2str(decorative)}", file=sys.stderr)

        alt_text: str = ""
        if not decorative:
            # get alt text from text file
            # print(f"Set to {csv_rows[image_cnt][6]}")
            alt_text = csv_rows[object_cnt][6]

        # set alt text
        if debug:
            print(f"Set group to {alt_text}")
        set_alt_text(group_shape, alt_text)

        slide_object_cnt += 1
        object_cnt += 1

        # remove last one
        group_shape_list = group_shape_list[:-1]

    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:

        # get decorative
        decorative_pptx:bool = is_decorative(shape)
        decorative:bool = str2bool(csv_rows[object_cnt][8])

        # change decorative status
        if decorative_pptx != decorative:
            # set decorative status of image
            print(f"Side: {slide_cnt}, {shape.name}, can't set the docorative status to: {bool2str(decorative)}", file=sys.stderr)

        alt_text: str = ""
        if not decorative:
            # get alt text from text file
            alt_text = csv_rows[object_cnt][6]

        # set alt text
        set_alt_text(shape, alt_text)

        slide_object_cnt += 1
        object_cnt += 1

    elif shape.shape_type in [MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.FREEFORM, \
                              MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.IGX_GRAPHIC, MSO_SHAPE_TYPE.CANVAS, \
                              MSO_SHAPE_TYPE.MEDIA, MSO_SHAPE_TYPE.WEB_VIDEO, MSO_SHAPE_TYPE.DIAGRAM, \
                              MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT, MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT, \
                              MSO_SHAPE_TYPE.WEB_VIDEO, MSO_SHAPE_TYPE.LINKED_OLE_OBJECT, \
                              MSO_SHAPE_TYPE.CALLOUT]:

        # get decorative
        decorative_pptx:bool = is_decorative(shape)
        decorative:bool = str2bool(csv_rows[object_cnt][8])

        # change decorative status
        if decorative_pptx != decorative:
            # set decorative status of image
            print(f"Side: {slide_cnt}, {shape.name}, can't set the docorative status to: {bool2str(decorative)}", file=sys.stderr)

        alt_text: str = ""
        if not decorative:
            # get alt text from text file
            alt_text = csv_rows[object_cnt][6]

        # set alt text
        set_alt_text(shape, alt_text)

        slide_object_cnt += 1
        object_cnt += 1

    elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:

        decorative_pptx:bool = is_decorative(shape)
        decorative:bool = str2bool(csv_rows[object_cnt][8])

        # change decorative status
        if decorative_pptx != decorative:
            # set decorative status of image
            print(f"Side: {slide_cnt}, {shape.name}, can't set the docorative status to: {bool2str(decorative)}", file=sys.stderr)

        alt_text: str = ""
        if not decorative:
            # get alt text from text file
            alt_text = csv_rows[object_cnt][6]

        # set alt text
        set_alt_text(shape, alt_text)

    return group_shape_list, object_cnt, slide_object_cnt

def add_presenter_note(
        pptx_path: str,
        pptx: dict,
        settings: dict,
        verbose: bool = False
    ) -> bool:
    """ add presenter note to each slide """
    err:bool = False

    pptx_name:str = pathlib.Path(pptx_path).stem
    #pptx_extension:str = pathlib.Path(pptx_path).suffix
    dirname:str = os.path.dirname(pptx_path)

    slide_dir:str = os.path.join(dirname, pptx_name, "slides_png")
    current_slide:int = pptx["slide_cnt"]
    #t = num2str(pptx['pptx_nslides'], current_slide + 1)
    t = str(current_slide + 1)
    file_name:str = f"{pptx_name}-{t}.png"
    slide_image_file_path = os.path.join(slide_dir, file_name)

    if os.path.isfile(slide_image_file_path):
        alt_text, err = generate_description(slide_image_file_path, ".png", settings, for_notes=True, verbose=verbose)
        slide = pptx["current_slide"]
        slide.notes_slide.notes_text_frame.text = alt_text
        if verbose:
            print(f"Slide: {current_slide + 1}\t{alt_text}")
    else:
        print(f"Unable to access image file: {slide_image_file_path}", file=sys.stderr)
        err = True

    return err

def remove_presenter_notes(
        pptx_path:str,
        verbose:bool = False
    ) -> bool:
    """ remove all presenter notes """
    err:bool = False

    # get name, extension, folder from Powerpoint file
    pptx_name:str = pathlib.Path(pptx_path).stem
    pptx_extension:str = pathlib.Path(pptx_path).suffix
    dirname:str = os.path.dirname(pptx_path)

    # process powerpoint file
    print(f"Processing Powerpoint file: {pptx_path}")
    prs = Presentation(pptx_path)

    # Loop through slides
    for _, slide in enumerate(prs.slides):
        slide.notes_slide.notes_text_frame.text = ""

    new_pptx_file_name = os.path.join(dirname, f"{pptx_name}_notes_removed{pptx_extension}")
    prs.save(new_pptx_file_name)
    if verbose:
        print(f"Saved Powerpoint file with presenter notes removed to: '{new_pptx_file_name}'\n")

    return err

def export_presenter_notes(
        pptx_path: str,
        verbose: bool = False
    ) -> bool:
    """ export presenter notes """
    err:bool = False

    # get name and folder from Powerpoint file
    pptx_name = pathlib.Path(pptx_path).stem
    dirname = os.path.dirname(pptx_path)

    # output file
    notes_file_path = os.path.join(dirname, f"{pptx_name}_notes.txt")

    with open(notes_file_path, "w", encoding="utf-8") as out_file:
        if verbose:
            print(f"Processing Powerpoint file: {pptx_path}")
        prs = Presentation(pptx_path)

        # Loop through slides
        for slide_cnt, slide in enumerate(prs.slides):
            if slide.shapes.title:
                title = slide.shapes.title.text
            else:
                title = ""
            s = f"=== Slide {slide_cnt} - {title} ===\n\n{slide.notes_slide.notes_text_frame.text}\n\n"
            out_file.write(s)

    if not err and verbose:
        print(f"Exported presenter notes to file: '{notes_file_path}'")

    return err

def export_slides_to_images(
        pptx_path: str,
        verbose: bool = False
    ) -> bool:
    """ export slides to PNG, Windows ONLY and requires that Powerpoint is installed """

    err:bool = False
    dirname:str = os.path.dirname(pptx_path)
    pptx_name:str = pathlib.Path(pptx_path).stem

    prs:Presentation = Presentation(pptx_path)
    n_slides = len(prs.slides)

    # create folder to store images
    img_folder = os.path.join(dirname, pptx_name)
    if not os.path.isdir(img_folder):
        os.makedirs(img_folder)

    path_to_folder_to_save = os.path.join(dirname, pptx_name, "slides_png")
    if not os.path.isdir(path_to_folder_to_save):
        os.makedirs(path_to_folder_to_save)

    if platform.system() == "Windows":
        try:
            import comtypes.client

            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            powerpoint.Visible = 1

            abs_file_path = os.path.abspath(pptx_path)
            presentation = powerpoint.Presentations.Open(abs_file_path)

            abs_path_to_folder_to_save = os.path.abspath(path_to_folder_to_save)

            for i, slide in enumerate(presentation.Slides):
                p = num2str(n_slides, i + 1)
                file_path = os.path.join(abs_path_to_folder_to_save, f"{pptx_name}-{p}.png")
                slide.Export(file_path, "PNG")

            presentation.Close()
            powerpoint.Quit()
            if verbose:
                print(f"Slides saved as PNG images in folder: '{abs_path_to_folder_to_save}'")
        except ImportError as e:
            print(f"Unable to export slides: {str(e)}", file=sys.stderr)
            err = True

    elif platform.system() != "Windows":
        # export PPTX first to PDF
        #print("Exporting to PDF...")

        cmd:list[str] = ["soffice", "--headless", "--convert-to", "pdf", pptx_path, "--outdir", path_to_folder_to_save]
        path_to_cmd = shutil.which(cmd[0])
        if path_to_cmd is not None:
            r = subprocess.run(cmd, stdin=subprocess.PIPE, stdout=subprocess.PIPE, shell=False, check=True)
        else:
            print("Warning, LibreOffice not installed.", file=sys.stderr)

        # save each page as a separate file
        if verbose:
            print("Extracting pages from PDF...")

        pdf_file:str = os.path.join(path_to_folder_to_save, f"{pptx_name}.pdf")
        out_file_name:str = os.path.join(path_to_folder_to_save, f"{pptx_name}.pdf")

        cmd = ["qpdf", "--split-pages", pdf_file, out_file_name]
        path_to_cmd = shutil.which(cmd[0])
        if path_to_cmd is not None:
            r = subprocess.run(cmd, stdin=subprocess.PIPE, stdout=subprocess.PIPE, shell=False, check=True)
        else:
            print("Warning, qpdf not installed.", file=sys.stderr)

        # export from PDF to PNG
        if verbose:
            print("Converting each page to PNG...")
        if platform.system() == "Darwin":
            the_files = os.listdir(path_to_folder_to_save)
            for f in the_files:
                if not f.startswith(".") and f.startswith(f"{pptx_name}-"):
                    in_file = os.path.join(path_to_folder_to_save, f)
                    cmd = ["sips", "-s", "dpiWidth", "300", "-s", "dpiHeight", "300", "-s", "format", "png", in_file, "--out", path_to_folder_to_save]
                    path_to_cmd = shutil.which(cmd[0])

                    if path_to_cmd is not None:
                        r = subprocess.run(cmd, stdin=subprocess.PIPE, stdout=subprocess.PIPE, shell=False, check=True) #capture_output = True)
                    else:
                        print("Unable to find 'sips'", file=sys.stderr)
        else:
            # Linux
            the_files = os.listdir(path_to_folder_to_save)
            for f in the_files:
                if not f.startswith(".") and f.startswith(f"{pptx_name}-"):
                    in_file = os.path.join(path_to_folder_to_save, f"{f}[0]")
                    out_file = f"{pathlib.Path(in_file).stem}.png"
                    cmd = ["convert", in_file, "-density", "300", out_file]
                    path_to_cmd = shutil.which(cmd[0])

                    if path_to_cmd is not None:
                        r = subprocess.run(cmd, stdin=subprocess.PIPE, stdout=subprocess.PIPE, shell=False, check=True) #capture_output = True)
                    else:
                        print("Unable to find 'convert'", file=sys.stderr)
        # remove PDFs
        for f in the_files:
            if not f.startswith(".") and f.startswith(f"{pptx_name}-"):
                in_file = os.path.join(path_to_folder_to_save, f)
                pid = os.remove(in_file)
                if pid:
                    _ = os.wait()

        # remove exported PDF file
        #os.remove(pdf_file)
        if verbose:
            print(f"Slides saved as PNG images in folder: '{path_to_folder_to_save}'")
    else:
        err = True
        print("Unable to convert PPTX to images.", file = sys.stderr)

    #if not err:
    #    # generate pptx from slide images
    #    images_to_pptx(path_to_folder_to_save, f"{pptx_name}_slides.pptx")

    return err

def images_to_pptx(
        images_path: str,
        output_file: str = "presentation.pptx",
        verbose: bool = False,
        debug: bool=False
    ) -> None:
    """
    Create a PowerPoint presentation with an image on each slide
    """
    print("Generate pptx from images...")

    # Create a presentation object
    prs = Presentation()

    # Define slide width and height (in centimeters)
    slide_width_cm = prs.slide_width.cm
    slide_height_cm = prs.slide_height.cm

    # select only PNG files
    the_files = os.listdir(images_path)
    the_files = [f for f in the_files if f.endswith(".png")]
    if len(the_files) == 0:
        print("No images found!")
        return

    for f in the_files:
        file_path = os.path.join(images_path, f)

        # Open the image to get its size
        with Image.open(file_path) as img:
            width_px, height_px = img.size
            dpi = img.info.get('dpi', (96, 96))  # Defaulting to 96 DPI if not provided

        # Calculate the image size in centimeters
        width_cm = width_px / dpi[0] * 2.54
        height_cm = height_px / dpi[1] * 2.54

        # Calculate the scaling factor
        scale_factor = min(slide_width_cm / width_cm, slide_height_cm / height_cm)

        # Add a slide with a blank layout
        slide_layout = prs.slide_layouts[5]  # 5 is the index for a blank slide
        slide = prs.slides.add_slide(slide_layout)

        # Add and resize the image to the slide
        img = slide.shapes.add_picture(file_path, Cm(0), Cm(0),
                                       width=Cm(width_cm * scale_factor),
                                       height=Cm(height_cm * scale_factor))

    # Save the presentation
    out = os.path.join(images_path, output_file)
    prs.save(out)

    if verbose:
        print(f"Saved pptx to '{out}'.")

def generate_description(
        image_file_path: str,
        extension: str, settings: dict,
        for_notes: bool = False,
        verbose: bool = False,
        debug: bool = False
    ) -> Tuple[str, bool]:
    """ generate image text description using MLLM/VL model """
    err:bool = False
    alt_text:str = ""
    model_str:str = settings["model"]

    if verbose:
        if for_notes:
            print("Generating presenter notes...")
        else:
            print("Generating alt text...")

    if settings["use_ollama"]:
        alt_text, err = use_ollama(image_file_path, settings, for_notes, verbose, debug)
    elif settings["use_mlx_vlm"]:
        alt_text, err = use_mlx_vlm(image_file_path, settings, for_notes, verbose, debug)
    else:
        if model_str == "kosmos-2":
            alt_text, err = kosmos2(image_file_path, settings, for_notes, verbose)
        elif model_str == "openclip":
            alt_text, err = openclip(image_file_path, settings, verbose)
        elif model_str == "qwen-vl":
            alt_text, err = qwen_vl(image_file_path, settings, for_notes, verbose)
        elif model_str == "cogvlm" or model_str == "cogvlm2":
            alt_text, err = cog_vlm(image_file_path, settings, for_notes, verbose)
        elif model_str == "phi3-vision":
            alt_text, err = phi3_vision(image_file_path, settings, for_notes, verbose)
        elif model_str == "gpt-4o" or model_str == "gpt-4-turbo":
            alt_text, err = use_openai(image_file_path, extension, settings, for_notes, verbose, debug)
        else:
            print(f"Unknown model: {model_str}", file = sys.stderr)

    # remove space at the end
    alt_text = alt_text.rstrip()

    return alt_text.lstrip(), err
