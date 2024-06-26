""" process.py """

from typing import Generator, List
import os
import sys
import argparse
import platform
import pathlib
import torch
import pandas as pd
from pytictoc import TicToc
from pptx import Presentation
from .powerpoint import process_shape, process_shapes_from_file, export_slides_to_images
from .powerpoint import add_presenter_note, remove_presenter_notes, export_presenter_notes
from .powerpoint import get_slide_img_path, generate_description
from .models import show_openclip_models, init_model

def process_images_from_pptx(
        file_path: str,
        settings: dict,
        debug: bool = False,
        verbose: bool = False
    ) -> Generator[int, None, None]:
    """
    Loop through images in the slides of a Powerpint file and set image description based 
    on image description from Kosmos-2, OpenCLIP, LLaVA, etc.
    """
    err: bool = False

    # get name, extension, folder from Powerpoint file
    pptx_name: str = pathlib.Path(file_path).stem
    pptx_extension: str = pathlib.Path(file_path).suffix
    dirname: str = os.path.dirname(file_path)

    report: bool = settings["report"]

    # create folder to store images
    img_folder:str = ""
    if not report:
        img_folder = os.path.join(dirname, pptx_name)
        if not os.path.isdir(img_folder):
            os.makedirs(img_folder)

    # Initialize presentation object
    if verbose:
        print(f"Reading PowerPoint file: '{file_path}'")
    prs:Presentation = Presentation(file_path)

    model_str:str = settings['model']

    # set output file name
    out_file_name:str = ""

    if not report:
        # generate alt text
        if settings["use_mlx_vlm"]:
            out_file_name = os.path.join(dirname, f"{pptx_name}_{model_str.replace('/', '_')}.json")
        else:
            out_file_name = os.path.join(dirname, f"{pptx_name}_{model_str.replace(':', '_')}.json")
    elif report:
        # just report
        out_file_name = os.path.join(dirname, f"{pptx_name}.json")

    pptx_nslides:int = len(prs.slides)

    # download and/or set up model
    if not report:
        err = init_model(settings)
        print()
        if err:
            print("Unable to init model.", file=sys.stderr)
            return err

    pptx:dict = {
        'group_shape_list': None,   # the group shape
        'image_list': None,    # list of images in the group shape
        'object_list': None,   # list of objects (except text boxes)
        'text_list': None,     # list of the text of text boxes in a shape group
        'base_left': 0,        # base_left of group shape
        'base_top': 0,         # base_top of group shape
        'pptx_name': pptx_name,
        'pptx_extension': pptx_extension,
        'fout': None,         # fout of text file
        'img_folder': img_folder,
        'pptx_nslides': pptx_nslides,
        'current_slide': None,
        'slide_cnt': 0,
        'slide_image_cnt': 0
    }

    # create pandas dataframe to store alt-text data
    col_names = ['Model', 'File', 'Slide', 'ShapeID', 'ObjectName', 'ObjectType', 'PartOfGroup', 'Alt_Text', \
                 'LenAltText', 'PresenterNotes', 'Decorative', 'Complex', 'PictFilePath', 'Complex_Alt_Text']
    df = pd.DataFrame(columns=col_names)
    pptx['df'] = df
    image_cnt:int = 0

    # Loop through slides
    print(f"Number of slides to process: {len(prs.slides)}")
    slide_cnt = 0
    for slide in prs.slides:

        # skip hidden slides
        if slide._element.get("show") == '0':
            continue

        pptx["slide_cnt"] = slide_cnt
        pptx["current_slide"] = slide
        if verbose:
            print(f"---- Slide: {slide_cnt + 1} ----")

        # loop through shapes
        pptx["slide_image_cnt"] = 0
        for shape in slide.shapes:
            err = process_shape(shape, pptx, settings, verbose, debug)
            if err:
                break

        if settings["add_to_notes"] and (pptx["slide_image_cnt"] > 0 or (pptx["object_list"] is not None)):
            # only add presenter note if there is at least one image or object on the slide
            err = add_presenter_note(file_path, pptx, settings, verbose)
        elif settings["add_to_notes_all_slides"]:
            # add description of slide to notes
            err = add_presenter_note(file_path, pptx, settings, verbose)
        else:
            slide = pptx["current_slide"]
            # keep or remove current presenter note
            if not report and not settings['keep_presenter_notes']:
                if slide.notes_slide.notes_text_frame is not None:
                    slide.notes_slide.notes_text_frame.text = ""
                else:
                    print(f"Unable to set slide notes: {slide_cnt + 1}", file=sys.stderr)

            model_str = settings['model']
            pptx_extension = pptx['pptx_extension']
            alt_text = ""
            presenter_notes = ""
            #if slide.notes_slide.notes_text_frame is not None:
            if slide.has_notes_slide:
                if slide.notes_slide.notes_text_frame is not None:
                    presenter_notes = slide.notes_slide.notes_text_frame.text
            else:
                print(f"Unable to set presenter note of slide: {slide_cnt + 1}", file=sys.stderr)
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                if text_frame is not None:
                    text_frame.text = presenter_notes
                else:
                    #text_frame = _BaseSlide.
                    #text_frame.text = ""
                    print("text frame error", file=sys.stderr)

            slide_image_file_path = get_slide_img_path(file_path, pptx)
            # check if exists because when creating accessibility report image might not
            # yet have been created
            if not os.path.isfile(slide_image_file_path):
                slide_image_file_path = ""

            df = pptx['df']
            # default entry
            df.loc[len(df)] = [
                model_str,
                f"{pptx_name}{pptx_extension}",
                pptx["slide_cnt"] + 1,
                0,
                "Slide",
                "",
                "",
                alt_text,
                len(alt_text),
                presenter_notes,
                False,
                False,
                slide_image_file_path,
                ""
            ]
            pptx['df'] = df

            # if err break out slide loop
            if err:
                break

            # reset info
            pptx["group_shape_list"] = None
            pptx["image_list"] = None
            pptx["object_list"] = None
            pptx["text_list"] = None

            image_cnt += pptx["slide_image_cnt"]

        yield slide_cnt + 1

        slide_cnt += 1

    if not err:
        # write to json
        df = pptx['df']
        df.to_json(out_file_name, orient='records', lines=True)

        if verbose:
            print("---------------------")
            print()
            print(f"Powerpoint file contains {len(prs.slides)} slides and in total {image_cnt} images with alt text.")

        pptx_file:str = ""
        if not report:
            # Save new pptx file
            if settings["use_mlx_vlm"]:
                new_pptx_file_name = os.path.join(dirname, f"{pptx_name}_{model_str.replace('/', '_')}{pptx_extension}")
            else:
                new_pptx_file_name = os.path.join(dirname, f"{pptx_name}_{model_str.replace(':', '_')}{pptx_extension}")

            print(f"\nSaving Powerpoint file with new alt-text to '{new_pptx_file_name}'\n")

            prs.save(new_pptx_file_name)
            pptx_file = new_pptx_file_name
        else:
            pptx_file = file_path

        # accessibility report
        print("---- Accessibility report --------------------------------------------")

        report = accessibility_report(df, pptx_file, debug)
        for l in report:
            print(l)

        print("----------------------------------------------------------------------")

    return err

def accessibility_report(df: pd.DataFrame, pptx_name: str, debug:bool = False) -> List[str]:
    """
    Create accessibility report based on infomation in the text file generated
    """
    lines = []
    lines.append(f"PowerPoint file: '{pptx_name}'\n")

    empty_alt_txt: int = 0
    empty_slide_presenter_notes: int = 0
    alt_text_list: list = []
    img_cnt: int = 0
    img_decorative_cnt: int = 0
    slide_cnt: int = 0
    group_cnt: int = 0
    other_objects: int = 0

    for i in range(len(df)):
        row = df.iloc[i]
        
        if row['ObjectType'] == "Picture":
            if not row['Decorative']:  #str2bool(row[9]):
                if row['LenAltText'] == 0:
                    # not decorative
                    empty_alt_txt += 1
                
                img_cnt += 1

                # create list of alt text length
                alt_text_list.append(row['LenAltText'])
            else:
                img_decorative_cnt += 1
        elif row['ObjectType'] == "Group":
            group_cnt += 1
        elif row['ObjectType'] == "" and row['ObjectName'] == "Slide":
            slide_cnt += 1
            if len(row['PresenterNotes']) == 0:
                empty_slide_presenter_notes += 1
        else:
            other_objects += 1

    lines.append(f"Slides: {slide_cnt}")
    lines.append(f"Elements: {len(df)}\n")
    lines.append(f"Images: {img_cnt}")
    lines.append(f"Decorative Images: {img_decorative_cnt}")
    lines.append(f"Groups: {group_cnt}")
    lines.append(f"Other Object: {other_objects}\n")

    lines.append(f"Number of missing Alt Text for non-decorative Images: {empty_alt_txt}")
    if len(alt_text_list) > 0:
        lines.append(f"Min Alt Text length: {min(alt_text_list)}")
        lines.append(f"Max Alt Text length: {max(alt_text_list)}")

    lines.append(f"\nNumber of Slides without Presenter Notes: {empty_slide_presenter_notes}")

    return lines

def update_pptx_df(
        file_path: str,
        file_path_json_file: str,
        df = None,
        save_to_original: bool = False,
        verbose: bool = False,
        debug:bool = False
    ) -> bool:
    """
    Replace alt texts specified in a json file (e.g. generated by this script and edited to correct or improve)
    """
    err:bool = False

    # get name, extension, folder from Powerpoint file
    name:str = pathlib.Path(file_path).stem
    extension:str = pathlib.Path(file_path).suffix
    dirname:str = os.path.dirname(file_path)

    if df is None:
        try:
            df = pd.read_json(file_path_json_file, orient='records', lines=True)
        except FileNotFoundError:
            print(f"FileNotFoundError: '{file_path_json_file}'", file=sys.stderr)
            err = True
    
    if verbose:
        print(f"Processing Powerpoint file: {file_path}")

    prs = Presentation(file_path)

    object_cnt: int = 1
    slide_cnt: int = 1
    for slide in prs.slides:

        # skip hidden slides
        if slide._element.get("show") == '0':
            continue

        slide_object_cnt = 0
        for shape in slide.shapes:
            _, object_cnt, slide_object_cnt = process_shapes_from_file(
                                                shape,
                                                None,
                                                df,
                                                slide_cnt,
                                                slide_object_cnt,
                                                object_cnt,
                                                verbose,
                                                debug
                                            )

        row = df.loc[
            (df['Slide'] == slide_cnt) &
            (df['ObjectName'] == "Slide") &
            (df['ObjectType'] == "")
        ]
        if not row.empty:
            presenter_notes = row.at[row.index[0], 'PresenterNotes']
            if slide.notes_slide.notes_text_frame is not None:
                slide.notes_slide.notes_text_frame.text = presenter_notes
            else:
                print(f"Error, unable to set presenter note of slide: {slide_cnt}", file=sys.stderr)
        else:
            print(f"Error, Slide presenter note not found on slide: {slide_cnt}", file=sys.stderr)

        slide_cnt += 1

    if not err:
        if save_to_original:
            if debug:
                print(f"Save to {file_path}")
            prs.save(file_path)
        else:
            outfile:str = os.path.join(dirname, f"{name}_alt_text{extension}")
            if verbose:
                print(f"Saving Powerpoint file with new alt-text to: '{outfile}'")
            prs.save(outfile)

    return err

def generate_alt_text_image(
        image_file_path: str,
        settings: dict,
        verbose: bool = False,
        debug:bool = False
) -> bool:
    """ generate alt text and save to text file """
    name:str = pathlib.Path(image_file_path).stem
    extension:str = pathlib.Path(image_file_path).suffix
    dirname:str = os.path.dirname(image_file_path)

    if debug:
        print(f"Generating Alt Text for {image_file_path}")
    alt_text, err = generate_description(
        image_file_path,
        extension,
        settings,
        for_notes=False,
        verbose=verbose
    )

    if not err:
        txt_file_path = os.path.join(dirname, f"{name}.txt")

        with open(txt_file_path, 'w', encoding='utf-8') as file:
            file.write(alt_text)

        if debug:
            print(f"alt text: {alt_text}")
            print(f"Saving to file: {txt_file_path}")

    return err

# argv: List[str]
def process_pptx() -> int:
    """ main """
    err:bool = False
    t = TicToc()
    t.tic()

    parser = argparse.ArgumentParser(description='Add Alt-Text automatically to images and objects in PowerPoint file')
    parser.add_argument("--pptx", type=str, default="", help="PowerPoint file")
    parser.add_argument("--image", type=str, default="", help="Image file")
    parser.add_argument("--report", action='store_true', default=False, help="flag to generate alt text report")
    parser.add_argument("--model", type=str, default="", help="kosmos-2, openclip, llava, gpt-4o, gpt-4-turbo, cogvlm, cogvlm2")

    #OpenAI
    parser.add_argument("--openai_api_key", type=str, default="", help="OpenAI API Key")

    # Ollama
    parser.add_argument("--use_ollama", action='store_true', default=False, help="use Ollama server")
    parser.add_argument("--server", type=str, default="http://localhost", help="Ollama server URL, default=http://localhost")
    parser.add_argument("--port", type=str, default="11434", help="Ollama server port, default=11434")

    # MLX-VLM
    parser.add_argument("--use_mlx_vlm", action='store_true', default=False, help="use MLX-VLM server")

    # OpenCLIP
    parser.add_argument("--show_openclip_models", action='store_true', default=False, help="show OpenCLIP models and pretrained models")
    parser.add_argument("--openclip_model", type=str, default="coca_ViT-L-14", help="OpenCLIP model")
    parser.add_argument("--openclip_pretrained", type=str, default="mscoco_finetuned_laion2B-s13B-b90k", help="OpenCLIP pretrained model")
    #
    parser.add_argument("--resize", type=str, default="500", help="resize image to same width and height in pixels, default:500, use 0 to disable resize")
    #
    parser.add_argument("--prompt", type=str, default="", help="custom prompt")
    parser.add_argument("--prompt_notes", type=str, default="", help="custom prompt for presenter notes")
    #
    #parser.add_argument("--save", action='store_true', default=False, help="flag to save powerpoint file with updated alt texts")
    #
    parser.add_argument("--keep_presenter_notes", action='store_true', default=False, help="replace or add to existing")
    #
    parser.add_argument("--replace", type=str, default="", help="replace alt texts in pptx with those specified in json file")
    parser.add_argument("--remove_presenter_notes", action='store_true', default=False, help="remove all presenter notes from powerpoint file")
    parser.add_argument("--export_presenter_notes", action='store_true', default=False, help="export presenter notes")
    parser.add_argument("--export_slides", action='store_true', default=False, help="export pptx slides to png images")
    parser.add_argument("--complex_alt_text", type=str, default="", help="alt text to use for complex images")
    #
    parser.add_argument("--add_to_notes", action='store_true', default=False, help="add slide description to slide notes when images are present")
    parser.add_argument("--add_to_notes_all_slides", action='store_true', default=False, help="add description to each slide notes")
    #
    parser.add_argument("--verbose", action='store_true', default=False, help="turn on verbose")
    parser.add_argument("--debug", action='store_true', default=False, help="flag for debugging")

    args = parser.parse_args()

    prompt:str = args.prompt
    model_str:str = args.model.lower()

    if args.show_openclip_models:
        show_openclip_models()
        return int(err)

    # set default prompt
    if model_str == "gpt-4-turbo" or model_str == "gpt-4o":
        if args.prompt == "":
            prompt = "Create a caption. Your response should be one or two sentences."
            #prompt = "Describe the image using one or two sentences. Do not mention the word 'image'."
    elif model_str == "kosmos-2":
        if args.prompt == "":
            prompt = "<grounding>An image of"
            #prompt = "<grounding>Describe this image:"
    elif model_str == "qwen-vl":
        if args.prompt == "":
            prompt = "Create a caption. Your response should be one or two sentences."
            #prompt = "Describe the image using one or two sentences."
    elif model_str == "cogvlm" or model_str == "cogvlm2":
        if args.prompt == "":
            prompt = "Create a caption. Your response should be one or two sentences."
            #prompt = "Describe the image using one or two sentences."
    elif args.use_ollama:
        if args.prompt == "":
            prompt = "Create a caption. Your response should be one or two sentences."
            #prompt = "You are an expert at understanding images and graphs. Answer concisely for someone who is visually impaired. Create a caption. Your response should be one or two sentences."
    elif args.use_mlx_vlm:
        if platform.system() != "Darwin":
            print("MLX-VLM is only available on MacOS.")
            return int(err)

        if args.prompt == "":
            prompt = "Create a caption. Your response should be one or two sentences."
    else:
        if args.prompt == "":
            prompt = "Create a caption. Your response should be one or two sentences."

    if args.prompt_notes == "":
        prompt_presenter_notes = "Describe the image in a very short paragraph without using bullet points or newlines for someone who is visually impaired. Start the desciption with 'This slide'"
    else:
        prompt_presenter_notes = args.prompt_notes

    # Read PowerPoint file and list images
    powerpoint_file_name = args.pptx
    image_file_name = args.image
    if powerpoint_file_name != "" and image_file_name == "":
        if not os.path.isfile(powerpoint_file_name):
            print(f"Error: PowerPoint file {powerpoint_file_name} not found", file=sys.stderr)
            err = True
    elif powerpoint_file_name == "" and image_file_name == "":
        if not os.path.isfile(powerpoint_file_name):
            print(f"Error: Image file '{image_file_name}' not found.", file=sys.stderr)
            err = True

    if not err:
        settings:dict = {
            "report": args.report,
            "model": model_str,
            "openai_api_key": args.openai_api_key,
            "kosmos2_model": None,
            "kosmos2_pretrained": None,
            "openclip_model_name": args.openclip_model,
            "openclip_pretrained": args.openclip_pretrained,
            "openclip-model": None,
            "openclip-preprocess": None,
            "qwen-vl-model": None,
            "qwen-vl-tokenizer": None,
            "cogvlm-model": None,
            "cogvlm-tokenizer": None,
            "use_ollama": args.use_ollama,
            "phi3-vision-model": None,
            "phi3-vision-tokenizer": None,
            "use_mlx_vlm": args.use_mlx_vlm,
            "ollama_url": f"{args.server}:{args.port}",
            "cuda_available": torch.cuda.is_available(),
            "mps_available": torch.backends.mps.is_available(),
            "prompt": prompt,
            "prompt_notes": prompt_presenter_notes,
            "img_size": int(args.resize),
            "keep_presenter_notes": args.keep_presenter_notes,
            "add_to_notes": args.add_to_notes,
            "add_to_notes_all_slides": args.add_to_notes_all_slides,
            "complex_alt_text": args.complex_alt_text
        }

        if image_file_name != "" and powerpoint_file_name == "":
            err = generate_alt_text_image(image_file_name, settings, args.verbose, args.debug)
        elif args.replace != "":
            # file with alt text provided
            err = update_pptx_df(powerpoint_file_name, args.replace, None, False, args.verbose, args.debug)
        elif args.remove_presenter_notes:
            err = remove_presenter_notes(powerpoint_file_name, args.verbose)
        elif args.export_presenter_notes:
            err = export_presenter_notes(powerpoint_file_name, args.verbose)
        elif args.export_slides:
            err = export_slides_to_images(powerpoint_file_name, args.verbose)
        else:
            if args.add_to_notes or args.add_to_notes_all_slides:
                print(f"Model: {model_str}")
                print(f"Presenter notes prompt: '{prompt_presenter_notes}'")
                
                print("Exporting slides to images...")
                t_export = TicToc()
                t_export.tic()
                err = export_slides_to_images(powerpoint_file_name, args.verbose)
                t_export.toc("Exporting slides took: ")

            # add alt-text
            if not err:
                for result in process_images_from_pptx(powerpoint_file_name, settings, args.debug, args.verbose):
                    print(f"Completed slide: {result}")

    print()
    elapsed = t.tocvalue()
    print(f"auto_alt_text took: {elapsed:.0f} seconds")

    return int(err)

if __name__ == "__main__":
    sys.exit(process_pptx())
