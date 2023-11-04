"""
Generate Alt Text for each picture in powerpoint files using different models (e.g. Kosmos-2, OpenCLIP, LLaVA)
"""

import os
import sys
import argparse
import requests
import base64
import csv
from pptx.oxml.ns import _nsmap
from typing import List
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape
import open_clip
import torch
from transformers import AutoProcessor, AutoModelForVision2Seq
from PIL import Image

def check_server_is_running(url: str) -> bool:
    try:
        response = requests.get(url)
        if response.status_code == 200:
            return True
    except requests.ConnectionError:
        return False
    return False

def num2str(the_max: int, n:int) -> str:
    if the_max > 99:
        if n < 100:
            if n < 10:
                return f"00{str(n)}"
            else:
                return f"0{str(n)}"
        else:
            return f"{str(n)}"
    else:
        if n < 10:
            return f"0{str(n)}"
        else:
            return f"{str(n)}"
        
def bool_value(s: str) -> bool:
    assert(s is not None and len(s) > 0)
    return s.lower() == "true"

def bool_to_string(b: bool) -> str:
    return "True" if b else "False"

# https://github.com/scanny/python-pptx/pull/512
def shape_get_alt_text(shape: BaseShape) -> str:
    """ Alt-text defined in shape's `descr` attribute, or '' if not present. """
    return shape._element._nvXxPr.cNvPr.attrib.get("descr", "")

def shape_set_alt_text(shape: BaseShape, alt_text: str):
    """ Set alt-text in shape """
    shape._element._nvXxPr.cNvPr.attrib["descr"] = alt_text

# https://stackoverflow.com/questions/63802783/check-if-image-is-decorative-in-powerpoint-using-python-pptx
def isDecorative(shape):
    # <adec:decorative xmlns:adec="http://schemas.microsoft.com/office/drawing/2017/decorative" val="1"/>
    _nsmap["adec"] = "http://schemas.microsoft.com/office/drawing/2017/decorative"
    cNvPr = shape._element._nvXxPr.cNvPr
    adec_decoratives = cNvPr.xpath(".//adec:decorative[@val='1']")
    if adec_decoratives:
        return True
    else:
        return False

def process_images_from_pptx(file_path: str, generate: bool, settings: dict, savePP: bool, DEBUG: bool = False) -> bool:
    """ 
    Loop through images in the slides of a Powerpint file and set image description based 
    on image description from Kosmos-2, OpenCLIP, or LLaVA
    """
    err: bool = False

    # get name, extension, folder from Powerpoint file
    file_name:str = os.path.basename(file_path)    
    name:str = file_name.split(".")[0]
    extension:str = file_name.split(".")[1]
    dirname:str = os.path.dirname(file_path)

    # Initialize presentation object
    print(f"Reading '{file_path}'")
    prs = Presentation(file_path)
    
    model_type:str = settings['model_type']

    # set output file name
    out_file_name:str = ""
    if model_type != "" and generate:
        out_file_name = os.path.join(dirname, f"{name}_{model_type}.txt")
    else:
        out_file_name = os.path.join(dirname, f"{name}.txt")
    
    nr_slides = len(prs.slides)

    # download and/or set up model
    if generate:
        err = init_model(settings)
        if err:
            print("Unable to init model.")
            return err

    # open file for writing
    fout = open(out_file_name, "w")

    # write header
    if model_type != "" and generate:
        fout.write(f"Model\tFile\tSlide\tPicture\tAlt_Text\tDecorative\tPict_File{os.linesep}")
    else:
        fout.write(f"File\tSlide\tPicture\tAlt_Text\Decorative\tPict_File{os.linesep}")

    # Loop through slides
    slide_cnt:int = 1
    image_cnt:int = 1
    image_file_path:str
    decorative:bool
    stored_alt_text:str
    for slide in prs.slides:
        # loop through shapes
        slide_image_cnt = 1
        for shape in slide.shapes:
            # Check if the shape has a picture
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE: 
                image_file_path = ""
                decorative = isDecorative(shape)

                # only generate alt text when generate options is True and decorative is False
                if generate and not decorative:
                    err, image_file_path = set_alt_text(shape, slide_cnt, nr_slides, slide_image_cnt, settings, DEBUG)
                
                # report alt text
                if not err:
                    stored_alt_text = shape_get_alt_text(shape)
                    feedback = f"Slide: {slide_cnt}, Picture: '{shape.name}', alt_text: '{stored_alt_text}', decorative: {bool_to_string(decorative)}"
                    print(feedback)

                    if model_type == "":
                        fout.write(f"{name}.{extension}\t{slide_cnt}\t{shape.name}\t{stored_alt_text}\t{bool_to_string(decorative)}\t{image_file_path}" + os.linesep)
                    else:
                        fout.write(f"{model_type}\t{name}.{extension}\t{slide_cnt}\t{shape.name}\t{stored_alt_text}\t{bool_to_string(decorative)}\t{image_file_path}" + os.linesep)

                    slide_image_cnt += 1
                    image_cnt += 1

        slide_cnt += 1

    # close output file
    fout.close()

    slide_cnt -= 1
    image_cnt -= 1
    print(f"Powerpoint file contains {slide_cnt} slides and {image_cnt} images.")

    if generate and savePP:
        # Save file
        outfile:str = os.path.join(dirname, f"{name}_alt_text.{extension}")
        print(f"Saving Powerpoint file with new alt-text to {outfile}")
        prs.save(outfile)

    return err

def init_model(settings: dict) -> bool:
    err: bool = False
    model_type:str = settings["model_type"]

    if model_type == "kosmos-2":
        # Kosmos-2 model
        model_name:str = "microsoft/kosmos-2-patch14-224"
        print(f"Kosmos-2 model: '{model_name}'")
        settings["kosmos2-model"] = AutoModelForVision2Seq.from_pretrained(model_name)
        settings["kosmos2-processor"] = AutoProcessor.from_pretrained(model_name)
    elif model_type == "openclip":
        # OpenCLIP
        print(f"OpenCLIP model: '{settings['openclip_model_name']}'\npretrained: '{settings['openclip_pretrained']}'")
        model, _, transform = open_clip.create_model_and_transforms(
            model_name=settings["openclip_model_name"],
            pretrained=settings["openclip_pretrained"]
        )
        settings["openclip-model"] = model
        settings["openclip-transform"] = transform
    elif model_type == "llava":
        # LLaVA
        server_url = settings["llava_url"]
        if check_server_is_running(server_url):
            server_url = f"{server_url}/completion"
            print(f"LLaVA server: '{server_url}'")
        else:
            print(f"Unable to access server at '{server_url}'.")
            err = True
    else:
        print(f"Unknown model: '{model_type}'")
        err = True

    return err

def set_alt_text(shape: BaseShape, slide_cnt: int, max_slides: int, image_cnt: int, settings: dict, DEBUG: bool) -> bool:
    err: bool = False
    
    # get image
    try:
        image_stream = shape.image.blob
        extension:str = shape.image.ext
    except Exception as e:
        print(f"Exception {str(e)}")
        return True
    
    image_file_name = f"s{num2str(max_slides, slide_cnt)}p{num2str(99, image_cnt)}_{shape.name}"
    image_file_path = os.path.join("tmp", image_file_name)
    image_file_path = os.path.join("tmp", f"{image_file_name}.{extension}")
    print(f"Saving and processing image: '{image_file_path}'...")

    # save image
    with open(image_file_path, "wb") as f:
        f.write(image_stream)

    alt_text: str = generate(image_file_path, settings)

    if DEBUG:
        print(f"Len: {len(alt_text)}, Content: {alt_text}")

    if len(alt_text) > 0:
        image_description = alt_text
        shape_set_alt_text(shape, image_description)
    else:
        print("No content.")

    return err, image_file_path

def generate(image_file_path: str, settings: dict, DEBUG:bool=False) -> str:
    alt_text: str = ""

    if settings["model_type"] == "kosmos-2":
        processor = settings["kosmos2-processor"]
        model = settings["kosmos2-model"]

        # read image
        im = Image.open(image_file_path)
        prompt = "<grounding>An image of"
        inputs = processor(text=prompt, images=im, return_tensors="pt")

        generated_ids = model.generate(
            pixel_values=inputs["pixel_values"],
            input_ids=inputs["input_ids"],
            attention_mask=inputs["attention_mask"],
            image_embeds=None,
            image_embeds_position_mask=inputs["image_embeds_position_mask"],
            use_cache=True,
            max_new_tokens=128,
        )
        generated_text = processor.batch_decode(generated_ids, skip_special_tokens=True)[0]

        # Specify `cleanup_and_extract=False` in order to see the raw model generation.
        #processed_text = processor.post_process_generation(generated_text, cleanup_and_extract=True)

        processed_text, entities = processor.post_process_generation(generated_text)

        # get picture description and remove trailing spaces
        alt_text = processed_text

    elif settings["model_type"] == "openclip":
        model = settings["openclip-model"]
        transform = settings["openclip-transform"]

        # read image
        im = Image.open(image_file_path).convert("RGB")
        im = transform(im).unsqueeze(0)

        # use OpenCLIP model to create label
        with torch.no_grad(), torch.cuda.amp.autocast():
            generated = model.generate(im)

        # get picture description and remove trailing spaces
        alt_text = open_clip.decode(generated[0]).split("<end_of_text>")[0].replace("<start_of_text>", "").strip()

        # remove space before '.'
        alt_text = alt_text.replace(' .', '.')
    elif settings["model_type"] == "llava":
        server_url = settings["llava_url"]
        server_url = f"{server_url}/completion"
        prompt = settings["llava_prompt"]

        # read image
        with open(image_file_path, 'rb') as img_file:
            img_byte_arr = img_file.read()

        # encode in base64
        img_base64 = base64.b64encode(img_byte_arr).decode('utf-8')

        # Use LLaVa to get image descriptions
        header = {"Content-Type": "application/json"}
        data = {
            "image_data": [{"data": img_base64, "id": 1}],
            "prompt": f"USER:[img-1] {prompt}\nASSISTANT:",
            "n_predict": 123,
            "temp": 0.1
        }
        try:
            response = requests.post(server_url, headers=header, json=data)
            response_data = response.json()

            if DEBUG:
                print(response_data)
                print()
        except Exception as e:
            print(f"Error: {str(e)}")
        else:
            # get picture description and remove trailing spaces
            alt_text = response_data.get('content', '').strip()

            # remove returns
            alt_text = alt_text.replace('\r', '')    

    return alt_text    

def add_alt_text_from_file(file_path: str, file_path_txt_file: str) -> bool:
    """
    Add alt text specified in a text file (e.g. generated by this script and edited to correct or improve)
    Text file should have a header and the same columns as the output files generated by this script
    """
    err: bool = False

    # Check if text file is exists
    if not os.path.isfile(file_path_txt_file):
        print(f"Unable to access file: {file_path_txt_file}")
        return False
    
    # get name, extension, folder from Powerpoint file
    file_name:str = os.path.basename(file_path)    
    name:str = file_name.split(".")[0]
    extension:str = file_name.split(".")[1]
    dirname:str = os.path.dirname(file_path)

    # process txt file
    print(f"Reading: {file_path_txt_file}...")
    csv_rows = []
    with open(file_path_txt_file, "r") as file:
        # assume tab delimited file
        csv_reader = csv.reader(file, delimiter="\t")

        # skip header
        next(csv_reader)

        for row in csv_reader:
            csv_rows.append(row)

    # process powerpoint file
    print(f"Processing Powerpoint file: {file_path}")
    prs = Presentation(file_path)

    # Loop through slides
    slide_cnt:int = 1
    image_cnt:int = 1
    for slide in prs.slides:
        # loop through shapes
        slide_image_cnt = 1
        for shape in slide.shapes:
            # Check if the shape has a picture
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                decorative_pptx = isDecorative(shape)

                # get decorative
                decorative = bool_value(csv_rows[image_cnt - 1][5])

                # change decorative status
                if decorative_pptx != decorative:
                    # set decorative status of image
                    print(f"Side: {slide_cnt}, {shape.name}, can't set the docorative status to: {bool_to_string(decorative)}")

                if decorative:
                    # decorative image
                    alt_text = ""
                    #print(f"Image is decorative")
                else:
                    # get alt text from text file
                    alt_text = csv_rows[image_cnt - 1][4]
                    #print(f"Setting alt-text to: {alt_text}")

                # set alt text
                shape_set_alt_text(shape, alt_text)

                slide_image_cnt += 1
                image_cnt += 1

    if not err:
        # Save file        
        outfile:str = os.path.join(dirname, f"{name}_alt_text.{extension}")
        print(f"Saving Powerpoint file with new alt-text to {outfile}")
        prs.save(outfile)

    return err


def main(argv: List[str]) -> int:
    err: bool = False
    
    default_llava_prompt:str = "Describe the image, figure, diagram, chart, table, or graph using a maximum of 125 characters"

    parser = argparse.ArgumentParser(description='Add alt-text automatically to images in Powerpoint')
    parser.add_argument("file", type=str, help="Powerpoint file")
    parser.add_argument("--generate", action='store_true', default=False, help="flag to generate alt-text to images")
    parser.add_argument("--type", type=str, default="", help="Model type: openclip, llava server, kosmos, gpt4")
    # LLaVA
    parser.add_argument("--prompt", type=str, default=default_llava_prompt, help="LLaVA prompt")
    parser.add_argument("--server", type=str, default="http://localhost", help="LLaVA server URL")
    parser.add_argument("--port", type=str, default="8007", help="LLaVA server port")
    # OpenCLIP
    parser.add_argument("--model", type=str, default="coca_ViT-L-14", help="model name")
    parser.add_argument("--pretrained", type=str, default="mscoco_finetuned_laion2B-s13B-b90k", help="pretrained model")
    #
    parser.add_argument("--save", action='store_true', default=False, help="Save powerpoint file")
    parser.add_argument("--add_from_file", type=str, default="", help="Add alt text from specified file to powerpoint file")
    #
    parser.add_argument("--debug", action='store_true', default=False, help="debug")

    args = parser.parse_args()

    # Read PowerPoint file and list images
    powerpoint_file_name = args.file
    if not os.path.isfile(powerpoint_file_name):
        print(f"Error: File {powerpoint_file_name} not found.")
        err = True
    else:
        model_type:str = args.type.lower()
        settings = {
            "model_type": model_type,
            "kosmos2_model": None,
            "kosmos2_pretrained": None,
            "openclip_model_name": args.model,
            "openclip_pretrained": args.pretrained,
            "openclip-model": None,
            "openclip-transform": None,
            "llava_prompt": args.prompt,
            "llava_url": f"{args.server}:{args.port}"
        }
        if args.add_from_file != "":
            err = add_alt_text_from_file(powerpoint_file_name, args.add_from_file)
        else:
            err = process_images_from_pptx(powerpoint_file_name, args.generate, settings, args.save, args.debug)

    return(int(err))

if __name__ == "__main__":
    exit_code = main(sys.argv[1:])
    sys.exit(exit_code)
