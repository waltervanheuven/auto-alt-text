"""
Generate Alt Text for each picture in a powerpoint file using MLLM and V-L pre-trained models
"""

from typing import Tuple
import os
import sys
import io
import subprocess
import shutil
import argparse
import base64
import platform
import csv
import json
import re
import pathlib
import requests
from urllib3.exceptions import HTTPError
from PIL import Image
import psutil
import open_clip
import torch
#import ollama
from transformers import AutoProcessor, AutoModelForVision2Seq, AutoModelForCausalLM, AutoTokenizer, LlamaTokenizer
#from transformers import LlavaNextProcessor, LlavaNextForConditionalGeneration
from transformers.generation import GenerationConfig
from pptx import Presentation
from pptx.util import Cm
from pptx.oxml.ns import _nsmap
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape

def check_server_is_running(url: str) -> bool:
    """ URL accessible? """    
    status:bool = False
    try:
        response = requests.get(url, timeout=10)
        if response.status_code == 200:
            status = True
    except requests.exceptions.Timeout:
        print("Timeout exception")
    except requests.exceptions.RequestException as e:
        print(f"Exception: {str(e)}")
    return status

def num2str(the_max: int, n:int) -> str:
    """ convert number to string with trailing zeros """
    s:str = f"{str(n)}"
    if the_max > 99:
        if n < 100:
            if n < 10:
                s = f"00{str(n)}"
            else:
                s = f"0{str(n)}"
    elif n < 10:
        s = f"0{str(n)}"
    return s

def str2bool(s: str) -> bool:
    """ convert str True or False to bool """
    assert(s is not None and len(s) > 0)
    return s.lower() == "true"

def bool2str(b: bool) -> str:
    """ convert bool to str """
    return "True" if b else "False"

# see https://github.com/scanny/python-pptx/pull/512
def get_alt_text(shape: BaseShape) -> str:
    """ Alt text is defined in shape's `descr` attribute, return this or '' if not present. """
    return shape._element._nvXxPr.cNvPr.attrib.get("descr", "")

def set_alt_text(shape: BaseShape, alt_text: str) -> None:
    """ Set alt text of shape """
    try:
        shape._element._nvXxPr.cNvPr.attrib["descr"] = alt_text
    except Exception as e:
        print(f"--> Unable to set alt_text: {shape.shape_type}, {shape.name}\n{str(e)}\nAlt_text: {alt_text}")

# see https://stackoverflow.com/questions/63802783/check-if-image-is-decorative-in-powerpoint-using-python-pptx
def is_decorative(shape) -> bool:
    """ check if image is decorative """
    # <adec:decorative xmlns:adec="http://schemas.microsoft.com/office/drawing/2017/decorative" val="1"/>
    _nsmap["adec"] = "http://schemas.microsoft.com/office/drawing/2017/decorative"
    cNvPr = shape._element._nvXxPr.cNvPr
    adec_decoratives = cNvPr.xpath(".//adec:decorative[@val='1']")
    return bool(adec_decoratives)

def process_images_from_pptx(file_path: str, settings: dict, debug: bool = False) -> bool:
    """
    Loop through images in the slides of a Powerpint file and set image description based 
    on image description from Kosmos-2, OpenCLIP, LLaVA, etc.
    """
    err:bool = False

    # get name, extension, folder from Powerpoint file
    pptx_name:str = pathlib.Path(file_path).stem
    pptx_extension:str = pathlib.Path(file_path).suffix
    dirname:str = os.path.dirname(file_path)

    report:bool = settings["report"]

    # create folder to store images
    img_folder:str = ""
    if not report:
        img_folder = os.path.join(dirname, pptx_name)
        if not os.path.isdir(img_folder):
            os.makedirs(img_folder)

    # Initialize presentation object
    print(f"Reading '{file_path}'")
    prs:Presentation = Presentation(file_path)

    model_str:str = settings['model']

    # set output file name
    out_file_name:str = ""

    if not report:
        # generate alt text
        out_file_name = os.path.join(dirname, f"{pptx_name}_{model_str.replace(":", "_")}.txt")
    elif report:
        # just report
        out_file_name = os.path.join(dirname, f"{pptx_name}.txt")

    pptx_nslides:int = len(prs.slides)

    # download and/or set up model
    if not report:
        err = init_model(settings)
        if err:
            print("Unable to init model.")
            return err
        else:
            print()

    pptx:dict = {
        'group_shape_list': None,   # the group shape
        'image_list': None,    # list of images in the group shape
        'object_list': None,   # list of objects (except text boxes)
        'text_list': None,     # list of the text of text boxes in a shape group
        'base_left': 0,        # base_left of group shape
        'base_top': 0,         # base_top of group shape
        'pptx_name': pptx_name,
        'pptx_extension': pptx_extension,
        'fout': None,         # fout of text file
        'img_folder': img_folder,
        'pptx_nslides': pptx_nslides,
        'current_slide': None,
        'slide_cnt': 0,
        'slide_image_cnt': 0
    }

    # open file for writing
    with open(out_file_name, "w", encoding="utf-8") as fout:
        # store fout
        pptx["fout"] = fout

        # write header
        fout.write("Model\tFile\tSlide\tObjectName\tObjectType\tPartOfGroup\tAlt_Text\tLenAltText\tDecorative\tPictFilePath\n")

        # total number of images in the pptx
        image_cnt:int = 0

        # Loop through slides
        slide_cnt:int = 0
        for slide_cnt, slide in enumerate(prs.slides):
            pptx["slide_cnt"] = slide_cnt
            pptx["current_slide"] = slide
            print(f"---- Slide: {slide_cnt + 1} ----")

            # loop through shapes
            pptx["slide_image_cnt"] = 0
            for shape in slide.shapes:
                err = process_shape(shape, pptx, settings, debug)
                if err:
                    break

            if settings["add_to_notes"] and (pptx["slide_image_cnt"] > 0 or (pptx["object_list"] is not None)):
                # only add presenter note if there is at least one image or object on the slide
                err = add_presenter_note(file_path, pptx, settings)

            # if err break out slide loop
            if err:
                break

            # reset info
            pptx["group_shape_list"] = None
            pptx["image_list"] = None
            pptx["object_list"] = None
            pptx["text_list"] = None

            image_cnt += pptx["slide_image_cnt"]

    if not err:
        print("---------------------")
        print()
        print(f"Powerpoint file contains {slide_cnt + 1} slides and in total {image_cnt} images with alt text.\n")

        pptx_file:str = ""
        if not report:
            # Save new pptx file
            new_pptx_file_name = os.path.join(dirname, f"{pptx_name}_{model_str.replace(":", "_")}{pptx_extension}")
            print(f"Saving Powerpoint file with new alt-text to '{new_pptx_file_name}'\n")
            prs.save(new_pptx_file_name)
            pptx_file = new_pptx_file_name
        else:
            pptx_file = file_path

        accessibility_report(out_file_name, pptx_file, debug)

    return err

def init_model(settings: dict) -> bool:
    """ download and init model for inference """
    err:bool = False
    model_str:str = settings["model"]
    prompt:str = settings["prompt"]
    model_name:str = ""
    model = None
    tokenizer = None

    if settings["use_ollama"]:
        print(f"Ollama server: {settings['ollama_url']}")
        
        # check if Ollama model is available on the server
        err, full_model_name = check_ollama_model_available(settings)

        if not err:
            settings['model'] = full_model_name
            print(f"Model: {settings['model']}")
            print(f"Prompt: '{prompt}'")
    elif model_str == "kosmos-2":
        # Kosmos-2 model
        model_name = "microsoft/kosmos-2-patch14-224"
        print(f"Kosmos-2 model: '{model_name}'")
        print(f"prompt: '{prompt}'")
        m = AutoModelForVision2Seq.from_pretrained(model_name)
        p = AutoProcessor.from_pretrained(model_name)
        if settings["cuda_available"]:
            print("Using CUDA.")
            m.cuda()
        settings["kosmos2-model"] = m
        settings["kosmos2-processor"] = p
    elif model_str == "openclip":
        # OpenCLIP
        print(f"OpenCLIP model: '{settings['openclip_model_name']}'\npretrained model: '{settings['openclip_pretrained']}'")

        #model, preprocess = open_clip.create_model_from_pretrained(settings["openclip_model_name"])
        #tokenizer = open_clip.get_tokenizer(settings["openclip_pretrained"])

        if settings["cuda_available"]:
            my_device = "cuda"
        else:
            print("Note that non-cuda devices are not support yet")
            my_device = "mps"
            err = True

        if not err:
            model, _, preprocess = open_clip.create_model_and_transforms(
                model_name=settings["openclip_model_name"],
                pretrained=settings["openclip_pretrained"],
                device=my_device,
                precision="fp16"
            )
            settings["openclip-model"] = model
            settings["openclip-preprocess"] = preprocess

    elif model_str == "qwen-vl":
        total_memory_bytes = psutil.virtual_memory().total
        total_memory_gb = total_memory_bytes / (1024**3)

        if settings["cuda_available"]:
            model_name = "Qwen/Qwen-VL-Chat"
            #model_name = "Qwen/Qwen-VL-Max"
            print(f"Qwen-VL model: '{model_name}'")
            print(f"prompt: '{prompt}'")
            print("Using CUDA.")

            tokenizer = AutoTokenizer.from_pretrained(model_name, trust_remote_code=True)
            model = AutoModelForCausalLM.from_pretrained(model_name, device_map="cuda", trust_remote_code=True).eval()
            model.generation_config = GenerationConfig.from_pretrained(model_name, trust_remote_code=True)
        elif settings['mps_available']:
            os.environ["ACCELERATE_USE_MPS_DEVICE"] = "True"
            model_name = "Qwen/Qwen-VL-Chat-Int4"
            #model_name = "Qwen/Qwen-VL"

            if total_memory_gb >= 32:
                print(f"Qwen-VL model: '{model_name}'")
                print(f"prompt: '{prompt}'")

                tokenizer = AutoTokenizer.from_pretrained(model_name, trust_remote_code=True)
                #model = AutoModelForCausalLM.from_pretrained(model_name, load_in_4bit=True, device_map="mps", trust_remote_code=True).eval()
                model = AutoModelForCausalLM.from_pretrained(model_name, device_map="mps", trust_remote_code=True).eval()
                model.generation_config = GenerationConfig.from_pretrained(model_name, trust_remote_code=True)
            else:
                print(f"Model '{model_name}' requires >= 32GB RAM.")
                err = True
        else:
            print(f"Model '{model_name}' requires a GPU with CUDA support")
            err = True

        settings["qwen-vl-model"] = model
        settings["qwen-vl-tokenizer"] = tokenizer
    elif model_str == "cogvlm":
        model_name = "THUDM/cogvlm-chat-hf"
        print(f"CogVLM model: '{model_name}'")
        print(f"prompt: '{prompt}'")

        if settings["cuda_available"]:
            print("Using CUDA.")

        tokenizer = LlamaTokenizer.from_pretrained('lmsys/vicuna-7b-v1.5')
        if settings["cuda_available"]:
            model = AutoModelForCausalLM.from_pretrained(
                model_name,
                load_in_4bit=True,
                #torch_dtype=torch.bfloat16,
                low_cpu_mem_usage=True,
                trust_remote_code=True
            ).to('cuda').eval()
        else:
            print("CogVLM requires a GPU with CUDA support.")
            err = True
            # if settings['mps_available']:
            #     os.environ["ACCELERATE_USE_MPS_DEVICE"] = "True"
            #     print("Activating mps device for accelerate.")

            # #     print("Not yet working on mps devices")
            # model = AutoModelForCausalLM.from_pretrained(
            #     model_name,
            #     load_in_4bit=True,
            #     #torch_dtype=torch.bfloat16,
            #     low_cpu_mem_usage=True,
            #     trust_remote_code=True,
            # ).to('mps').eval()
            # # else:
            #     print("Not yet working on cpu")
            #     model = AutoModelForCausalLM.from_pretrained(
            #         model_name,
            #         load_in_4bit=True,
            #         trust_remote_code=True,
            #     ).to('cpu').eval()

        settings["cogvlm-model"] = model
        settings["cogvlm-tokenizer"] = tokenizer

    elif model_str == "gpt-4v":
        print("GPT-4V")
        print(f"model: {settings['gpt4v_model']}")
        print(f"prompt: '{prompt}'")
    else:
        print(f"Unknown model: '{model_str}'")
        err = True

    return err

def process_shape(shape: BaseShape, pptx: dict, settings: dict, debug: bool) -> bool:
    """
    Recursive function to process shapes and shapes in groups on each slide
    """
    err: bool = False
    report:bool = settings["report"]

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        # keep a list of images as part of the group
        if pptx["group_shape_list"] is None:
            pptx["group_shape_list"] = [shape]
        else:
            group_shape_list = pptx["group_shape_list"]
            group_shape_list.append(shape)
            pptx["group_shape_list"] = group_shape_list

        if pptx["image_list"] is None:
            pptx["image_list"] = []
        if pptx["text_list"] is None:
            pptx["text_list"] = []
        if pptx["object_list"] is None:
            pptx["object_list"] = []

        # process shapes
        for embedded_shape in shape.shapes:
            err = process_shape(embedded_shape, pptx, settings, debug)
            if err:
                break

        if not err:
            # check if group is not part of other group
            group_shape_list:list[BaseShape] = pptx["group_shape_list"]
            part_of_group:str = "No"
            if len(group_shape_list) > 1:
                part_of_group = "Yes"
            elif len(group_shape_list) == 1:
                part_of_group = "No_TopLevel"

            # current group shape
            group_shape:BaseShape = get_current_group_shape(pptx)

            # image list
            image_list:list = pptx["image_list"]

            # group contains at least one image
            # if image_list is not None and len(image_list) > 0:
            #     new_img = combine_images_in_group(image_list, group_shape)
            #     img_folder = pptx["img_folder"]
            #     filename = os.path.join(img_folder, f'slide_{pptx["slide_cnt"]}_group.png')
            #     new_img.save(filename)

            alt_text:str = ""
            if not report:
                # combine text box content associated with group
                text_list:list = pptx["text_list"]
                for n, txt in enumerate(text_list):
                    # remove newlines
                    txt = replace_newline_with_space(txt)
                    if n == 0:
                        alt_text = txt
                    else:
                        alt_text = f"{alt_text} {txt}"

                if len(alt_text) > 0:
                    alt_text = f"{alt_text}. "

                # combine alt text to generate the alt text for the group

                if len(image_list) > 1:
                    alt_text = f"{alt_text}There are {len(image_list)} images:"
                for _, _, _, txt in image_list:
                    # remove newlines
                    txt = replace_newline_with_space(txt)
                    if len(alt_text) == 0:
                        alt_text = txt
                    else:
                        alt_text = f"{alt_text} {txt}"

                # set alt text of group shape
                set_alt_text(group_shape, alt_text)
            else:
                alt_text = get_alt_text(group_shape)

                # remove returns
                alt_text = replace_newline_with_space(alt_text)

            # get vars
            model_str:str = settings["model"]
            report:bool = settings["report"]
            image_file_path:str = ""
            pptx_name:str = pptx["pptx_name"]
            pptx_extension:str = pptx["pptx_extension"]
            slide_cnt:int = pptx["slide_cnt"]

            # get info from groupshape
            decorative:bool = is_decorative(group_shape)
            stored_alt_text:str = get_alt_text(group_shape)

            if decorative:
                print(f"Slide: {slide_cnt + 1}, Group: {group_shape.name}, alt_text: '{stored_alt_text}', decorative")
            else:
                print(f"Slide: {slide_cnt + 1}, Group: {group_shape.name}, alt_text: '{stored_alt_text}'")

            fout = pptx["fout"]
            fout.write(f"{model_str}\t{pptx_name}{pptx_extension}\t{slide_cnt + 1}\t{group_shape.name}\tGroup\t{part_of_group}\t{stored_alt_text}\t{len(stored_alt_text)}\t{bool2str(decorative)}\t{image_file_path}\n")

            # remove last one
            group_shape_list = pptx["group_shape_list"]
            pptx["group_shape_list"] = group_shape_list[:-1]

    elif shape.shape_type in [MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE]:
        # picture
        image_file_path:str = ""
        decorative:bool = is_decorative(shape)
        group_shape:BaseShape = get_current_group_shape(pptx)

        # only generate alt text when generate options is True and decorative is False
        if not decorative:
            err, image_file_path = process_shape_and_generate_alt_text(shape, pptx, settings, debug)

        if not err:
            part_of_group = "No"
            if group_shape is not None:
                part_of_group = "Yes"

            # report alt text
            if not err:
                slide_cnt:int = pptx["slide_cnt"]
                slide_image_cnt:int = pptx["slide_image_cnt"]

                stored_alt_text = get_alt_text(shape)
                if decorative:
                    print(f"Slide: {slide_cnt + 1}, Pict: {slide_image_cnt + 1}, alt_text: '{stored_alt_text}', decorative")
                else:
                    print(f"Slide: {slide_cnt + 1}, Pict: {slide_image_cnt + 1}, alt_text: '{stored_alt_text}'")

                model_str:str = settings["model"]
                pptx_name:str = pptx["pptx_name"]
                pptx_extension:str = pptx["pptx_extension"]
                fout = pptx["fout"]
                fout.write(f"{model_str}\t{pptx_name}{pptx_extension}\t{slide_cnt + 1}\t{shape.name}\tPicture\t{part_of_group}\t{stored_alt_text}\t{len(stored_alt_text)}\t{bool2str(decorative)}\t{image_file_path}\n")

                pptx["slide_image_cnt"] = slide_image_cnt + 1

    elif shape.shape_type in [MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.FREEFORM, \
                              MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.IGX_GRAPHIC, MSO_SHAPE_TYPE.CANVAS, \
                              MSO_SHAPE_TYPE.MEDIA, MSO_SHAPE_TYPE.WEB_VIDEO, MSO_SHAPE_TYPE.DIAGRAM, \
                              MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT, MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT, \
                              MSO_SHAPE_TYPE.WEB_VIDEO, MSO_SHAPE_TYPE.LINKED_OLE_OBJECT, \
                              MSO_SHAPE_TYPE.CALLOUT]:

        process_object(shape, pptx, settings, debug)

        object_list = pptx["object_list"]
        if object_list is None:
            object_list = []
            object_list.append(shape)
            pptx["object_list"] = object_list

    elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
        # TEXT_BOX is part of a group
        # For the Alt Text it would be useful to add this text to the Alt Text
        group_shape:BaseShape = get_current_group_shape(pptx)
        if group_shape is not None:
            text:str = shape.text_frame.text
            text_list = pptx["text_list"]
            if text_list is not None:
                text_list.append(text)
                pptx["text_list"] = text_list
    elif debug:
        print(f"=> OBJECT: {shape.name}, type: {shape.shape_type}")

    return err

def replace_newline_with_space(txt: str) -> str:
    """ replace newline with space and replace tab with space """
    s = " ".join(txt.splitlines())
    return s.replace("\t", " ")

def get_current_group_shape(pptx:dict) -> BaseShape:
    """ return group shape """
    group_shape_list:list[BaseShape] = pptx["group_shape_list"]
    if group_shape_list is not None and len(group_shape_list) > 0:
        return group_shape_list[-1]

    return None

def shape_type2str(type) -> str:
    """ return name of shape as string """
    s:str = ""

    if type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        s =  "Auto shape"
    elif type == MSO_SHAPE_TYPE.LINE:
        s =  "Line"
    elif type == MSO_SHAPE_TYPE.IGX_GRAPHIC:
        s =  "IgxGraphic"
    elif type == MSO_SHAPE_TYPE.CHART:
        s = "Chart"
    elif type == MSO_SHAPE_TYPE.FREEFORM:
        s = "FreeForm"
    elif type == MSO_SHAPE_TYPE.TEXT_BOX:
        s = "TextBox"
    elif type == MSO_SHAPE_TYPE.CANVAS:
        s = "Canvas"
    elif type == MSO_SHAPE_TYPE.MEDIA:
        s = "Media"
    elif type == MSO_SHAPE_TYPE.WEB_VIDEO:
        s = "WebVideo"
    elif type == MSO_SHAPE_TYPE.DIAGRAM:
        s = "Diagram"
    elif type == MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT:
        s = "Control object"
    elif type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
        s = "Embedded object"
    elif type == MSO_SHAPE_TYPE.WEB_VIDEO:
        s = "Web video"

    return s

def process_object(shape:BaseShape, pptx:dict, settings:dict, debug:bool = False) -> None:
    """ process """
    # only include if it is not part of a group
    # Powerpoint only reports an accessibility error for a missing group shape alt text
    image_file_path:str = ""
    decorative:bool = is_decorative(shape)
    report:bool = settings["report"]

    # include all text inside shape?
    include_all_paragraphs = True

    group_shape:BaseShape = get_current_group_shape(pptx)
    part_of_group:str = "No"
    if group_shape is not None:
        part_of_group = "Yes"

    alt_text:str = ""
    if not report:
        # Quick fix for alt text of shapes
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if len(shape.name) > 0:
                alt_text = f"{cleanup_name_object(shape.name)}"
            else:
                alt_text = ""
        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            if len(shape.name) > 0:
                if cleanup_name_object(shape.name.lower()) != "chart":
                    # avoid duplication chart
                    alt_text = f"{cleanup_name_object(shape.name)}"
                else:
                    alt_text = f"{cleanup_name_object(shape.name)}"
            else:
                alt_text = ""

            # add title of chart to alt_text
            the_chart = shape.chart
            if the_chart.has_title and len(the_chart.chart_title.text_frame.text.strip()) > 0:
                alt_text = f"{alt_text} entitled '{the_chart.chart_title.text_frame.text}'"

        elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
            if len(shape.name) > 0:
                alt_text = f"{cleanup_name_object(shape.name)}"
            else:
                alt_text = ""
        elif shape.shape_type == MSO_SHAPE_TYPE.CANVAS:
            if len(shape.name) > 0:
                alt_text = f"{cleanup_name_object(shape.name)}"
            else:
                alt_text = ""
        elif shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
            if len(shape.name) > 0:
                alt_text = f"{cleanup_name_object(shape.name)}"
            else:
                alt_text = ""
        elif shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
            if len(shape.name) > 0:
                alt_text = f"{cleanup_name_object(shape.name)}"
            else:
                alt_text = ""
        elif shape.shape_type == MSO_SHAPE_TYPE.WEB_VIDEO:
            if len(shape.name) > 0:
                alt_text = f"{cleanup_name_object(shape.name)}"
            else:
                alt_text = ""
        elif shape.shape_type == MSO_SHAPE_TYPE.DIAGRAM:
            if len(shape.name) > 0:
                alt_text = f"{shape.name}"
            else:
                alt_text = ""
        elif shape.shape_type == MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT:
            if len(shape.name) > 0:
                alt_text = f"{shape.name}"
            else:
                alt_text = ""
        elif shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
            if len(shape.name) > 0:
                alt_text = f"{shape.name}"
            else:
                alt_text = ""
        else:
            if len(shape.name) > 0:
                alt_text = f"{shape.name.lower()}"
            else:
                alt_text = ""

        # indicate the text inside the shape
        if shape.has_text_frame:
            if len(shape.text_frame.paragraphs) == 1 and shape.text_frame.paragraphs[0].text != "":
                alt_text = f"{alt_text} with inside the text: {remove_newlines(shape.text_frame.paragraphs[0].text).strip()}"
            else:
                if not include_all_paragraphs:
                    alt_text = f"{alt_text} with text inside."
                else:
                    first = True
                    for p in shape.text_frame.paragraphs:
                        if p.text != "":
                            if first:
                                alt_text = f"{alt_text} with inside the text:"
                                first = False

                            alt_text = f"{alt_text} {remove_newlines(p.text).strip()}"

        # make sure alt text ends with a final stop
        if not alt_text.endswith("."):
            alt_text = f"{alt_text}."

        set_alt_text(shape, alt_text)

        # if part of group store alt_text
        if group_shape is not None:
            text_list = pptx["text_list"]
            if text_list is not None:
                text_list.append(alt_text)
                pptx["text_list"] = text_list

    # report alt text
    slide_cnt:int = pptx["slide_cnt"]
    slide_image_cnt:int = pptx["slide_image_cnt"]

    stored_alt_text = get_alt_text(shape)
    if decorative:
        print(f"Slide: {slide_cnt + 1}, {shape_type2str(shape.shape_type)}: {slide_image_cnt + 1}, alt_text: '{stored_alt_text}', decorative")
    else:
        print(f"Slide: {slide_cnt + 1}, {shape_type2str(shape.shape_type)}: {slide_image_cnt + 1}, alt_text: '{stored_alt_text}'")

    model_str:str = settings["model"]
    pptx_name:str = pptx["pptx_name"]
    pptx_extension:str = pptx["pptx_extension"]
    fout = pptx["fout"]
    fout.write(f"{model_str}\t{pptx_name}{pptx_extension}\t{slide_cnt + 1}\t{shape.name}\t{shape_type2str(shape.shape_type)}\t{part_of_group}\t{stored_alt_text}\t{len(stored_alt_text)}\t{bool2str(decorative)}\t{image_file_path}\n")

def remove_newlines(txt:str) -> str:
    """ remove newlines and replace tabs with spaces """
    s = "".join(txt.splitlines())
    return s.replace("\t", " ")

def cleanup_name_object(txt:str) -> str:
    """
    check if alt shape name contains a number at the end 
    e.g. "oval 1", "oval 2" and remove the number
    """
    s:str = ""

    elements:list[str] = txt.lower().split()
    if len(elements) == 1:
        s = elements[0]
    else:
        last_word = elements[-1]
        try:
            _ = int(last_word)
            s = ' '.join(elements[:-1])
        except ValueError:
            s = txt

    return s

def combine_images_in_group(images, group_shape) -> Image:
    """ 
    Create new image based on shape

    TODO: Not yet working properly, image size is not correct
    """

    # EMU per Pixel estimate: not correct
    EMU_PER_PIXEL:int = int(914400 / 300)

    # Determine the size of the new image based on the group shape size
    new_img_width = int(group_shape.width / EMU_PER_PIXEL)
    new_img_height = int(group_shape.height / EMU_PER_PIXEL)
    new_img = Image.new('RGB', (new_img_width, new_img_height))

    # Paste each image into the new image at its relative position
    for image, left, top, _ in images:
        print(f"img: {image.width}, {image.height}, {left}, {top}")
        new_img.paste(image, (int(left / EMU_PER_PIXEL), int(top / EMU_PER_PIXEL)))

    return new_img

def process_shape_and_generate_alt_text(shape:BaseShape, pptx:dict, settings:dict, debug:bool=False) -> Tuple[bool, str]:
    """ 
    Save image associated with shape and generate alt text
    """
    err:bool = False
    image_file_path:str = ""

    # get image
    image_stream = None
    extension:str = ""
    if hasattr(shape, "image"):
        # get image, works with only with PNG, JPEG?
        image_stream = shape.image.blob
        extension = shape.image.ext
    else:
        # get image for other formats, e.g. TIFF
        # <Element {http://schemas.openxmlformats.org/presentationml/2006/main}pic at 0x15f2d6b20>
        try:
            slide_part = shape.part
            rId = shape._element.blip_rId
            image_part = slide_part.related_part(rId)
            image_stream = image_part.blob
            extension = image_part.partname.ext
        except AttributeError:
            slide_cnt:int = pptx["slide_cnt"] + 1
            print(f"Error, slide: {slide_cnt}, pict: '{shape.name}', unable to access image")
            #err = True

    if not err and image_stream is not None:
        report:bool = settings["report"]
        base_left:int = pptx["base_left"]
        base_top:int = pptx["base_top"]

        # determine file name
        pptx_nslides: int = pptx["pptx_nslides"]
        slide_image_cnt:int = pptx["slide_image_cnt"]
        slide_cnt:int = pptx["slide_cnt"]
        pptx_nslides:int = pptx["pptx_nslides"]
        img_folder:str = pptx["img_folder"]

        alt_text:str = ""
        if not report:
            image_file_name:str = f"s{num2str(pptx_nslides, slide_cnt + 1)}p{num2str(99, slide_image_cnt + 1)}"
            image_file_path = os.path.join(img_folder, f"{image_file_name}.{extension}")
            print(f"Saving image from pptx: '{image_file_path}'")

            # save image
            with open(image_file_path, "wb") as f:
                f.write(image_stream)

            alt_text, err = generate_description(image_file_path, extension, settings)
        else:
            alt_text = get_alt_text(shape)

            # remove returns
            alt_text = replace_newline_with_space(alt_text)

        if not err:
            # Keep image in case the image is part of a group
            if pptx["image_list"] is not None:
                # Calculate the position of the image relative to the group
                image_group_part = Image.open(io.BytesIO(image_stream))
                left = base_left + shape.left
                top = base_top + shape.top
                image_list = pptx["image_list"]
                image_list.append((image_group_part, left, top, alt_text))
                pptx["image_list"] = image_list

            if debug:
                print(f"Len: {len(alt_text)}, Content: {alt_text}")

            if len(alt_text) > 0:
                set_alt_text(shape, alt_text)

    return err, image_file_path

def generate_description(image_file_path: str, extension:str, settings: dict, for_notes:bool=False, debug:bool=False) -> Tuple[str, bool]:
    """ generate image text description using MLLM/VL model """
    err:bool = False
    alt_text:str = ""
    model_str:str = settings["model"]

    if for_notes:
        print("Generating presenter notes...")
    else:
        print("Generating alt text...")

    if settings["use_ollama"]:
        alt_text, err = use_ollama(image_file_path, extension, settings, for_notes, debug)
    else:
        if model_str == "kosmos-2":
            alt_text, err = kosmos2(image_file_path, settings, for_notes, debug)
        elif model_str == "openclip":
            alt_text, err = openclip(image_file_path, settings, for_notes, debug)
        elif model_str == "qwen-vl":
            alt_text, err = qwen_vl(image_file_path, settings, for_notes, debug)
        elif model_str == "cogvlm":
            alt_text, err = cog_vlm(image_file_path, settings, for_notes, debug)
        elif model_str == "gpt-4v":
            alt_text, err = gpt4v(image_file_path, extension, settings, for_notes, debug)
        else:
            print(f"Unknown model: {model_str}")

    # remove space at the end
    alt_text = alt_text.rstrip()

    return alt_text.lstrip(), err

def kosmos2(image_file_path: str, settings: dict, for_notes:bool=False, debug:bool=False) -> Tuple[str, bool]:
    """ get image description from Kosmos-2 """
    err:bool = False

    # check if readonly
    image_file_path, readonly, msg = check_readonly_formats(image_file_path, settings)
    if readonly:
        return msg, False

    with Image.open(image_file_path) as img:

        # resize image
        img = resize(img, settings)

        # prompt
        if for_notes:
            prompt = settings["prompt_notes"]
        else:
            prompt = settings["prompt"]

        processor:str = settings["kosmos2-processor"]
        model:str = settings["kosmos2-model"]

        inputs = processor(text=prompt, images=img, return_tensors="pt")
        if settings["cuda_available"]:
            generated_ids = model.generate(
                pixel_values=inputs["pixel_values"].cuda(),
                input_ids=inputs["input_ids"].cuda(),
                attention_mask=inputs["attention_mask"].cuda(),
                image_embeds=None,
                image_embeds_position_mask=inputs["image_embeds_position_mask"].cuda(),
                use_cache=True,
                max_new_tokens=256,
            )
        else:
            generated_ids = model.generate(
                pixel_values=inputs["pixel_values"],
                input_ids=inputs["input_ids"],
                attention_mask=inputs["attention_mask"],
                image_embeds=None,
                image_embeds_position_mask=inputs["image_embeds_position_mask"],
                use_cache=True,
                max_new_tokens=256,
            )

    generated_text = processor.batch_decode(generated_ids, skip_special_tokens=True)[0]

    # Specify `cleanup_and_extract=False` in order to see the raw model generation.
    #processed_text = processor.post_process_generation(generated_text, cleanup_and_extract=True)

    # processed_text, entities = processor.post_process_generation(generated_text)
    processed_text, _ = processor.post_process_generation(generated_text)

    # remove prompt
    p:str = re.sub('<[^<]+?>', '', prompt)
    processed_text = processed_text.replace(p.strip(), '')

    # capitalize
    alt_text:str = processed_text.strip().capitalize()

    return alt_text, err

def check_readonly_formats(image_file_path: str, settings: dict) -> Tuple[str, str, bool]:
    """
    Check if image format is WMF, WME, or PSD which can not be converted using the pillow library.

    Function converts WMF (vector format) to JPEG using LibreOffice.
    
    Other read only formats not yet tested. Conversion only tested on macOS.
    """
    readonly:bool = False
    msg:str = ""
    new_image_file_path = image_file_path

    with Image.open(image_file_path) as img:

        if img.format in ['WMF', 'WME']:
            msg = "A windows media format file."
            readonly = True
        elif img.format in ['PSD']:
            msg = "An Adobe Photoshop file."
            readonly = True

        if readonly and img.format in ['WMF']:
            err:bool = False

            # convert images to PNG
            dirname:str = os.path.dirname(image_file_path)
            basename:str = os.path.basename(image_file_path).split(".")[0]
            new_image_file_path = os.path.join(os.path.dirname(image_file_path), f"{basename}.png")

            print(f"Converting {img.format} to PNG...")
            try:
                # convert WMF to PNG using headless libreoffice
                if platform.system() != "Windows":
                    # convert using LibreOffice (headless)
                    cmd:list[str] = ["soffice", "--headless", "--convert-to", "png", image_file_path, "--outdir", dirname]
                    path_to_cmd = shutil.which(cmd[0])
                    if path_to_cmd is not None:
                        r = subprocess.run(cmd, stdin=subprocess.PIPE, stdout=subprocess.PIPE, shell=False, check=True)
                    else:
                        print("Warning, LibreOffice not installed.")
                elif platform.system() == "Windows":
                    # convert using magick
                    cmd:list[str] = ["magick", "convert", image_file_path, new_image_file_path]
                    path_to_cmd = shutil.which(cmd[0])
                    if path_to_cmd is not None:
                        r = subprocess.run(cmd, stdin=subprocess.PIPE, stdout=subprocess.PIPE, shell=False, check=True)
                    else:
                        print("Warning, ImageMagick not installed.")
            except subprocess.CalledProcessError as e:
                msg = f"soffice CalledProcessError: {str(e)}"
                err = True
            except subprocess.TimeoutExpired as e:
                msg = f"soffice TimeoutExpired: {str(e)}"
                err = True
            except OSError as e:
                msg = f"soffice OSError, file does not exist?: {str(e)}"
                err = True
            except Exception as e:
                msg = f"soffice exception: {str(e)}"
                err = True
            else:
                readonly = False

            if err:
                readonly = True
                new_image_file_path = image_file_path
                print(r.stderr)

    if readonly:
        print(f"Warning, unable to open '{img.format}' file. Replace image in powerpoint with PNG, TIFF, or JPEG version.")

    return new_image_file_path, readonly, msg

def resize(image:Image.Image, settings:dict) -> Image.Image:
    """ resize image """
    px:int = settings["img_size"]
    if px != 0:
        # only resize if img_size != 0
        if image.width > px or image.height > px:
            new_size = (min(px, image.width), min(px, image.height))
            print(f"Resize image from ({image.width}, {image.height}) to {new_size}")
            image = image.resize(new_size)

    return image

def openclip(image_file_path: str, settings: dict, for_notes:bool=False, debug:bool=False) -> Tuple[str, bool]:
    """ get image description from OpenCLIP """
    err:bool = False

    # check if readonly
    image_file_path, readonly, msg = check_readonly_formats(image_file_path, settings)
    if readonly:
        return msg, False

    with Image.open(image_file_path).convert('RGB') as img:
        # resize image
        img = resize(img, settings)

        preprocess = settings["openclip-preprocess"]
        img = preprocess(img).unsqueeze(0)

    # use OpenCLIP model to create label
    model = settings["openclip-model"]

    if settings["cuda_available"]:
        with torch.no_grad(), torch.cuda.amp.autocast():
            generated = model.generate(img)
    else:
        with torch.no_grad(): #, torch.autocast('mps'):
            generated = model.generate(img)

    # get picture description and remove trailing spaces
    alt_text = open_clip.decode(generated[0]).split("<end_of_text>")[0].replace("<start_of_text>", "").strip()

    # remove space before '.' and capitalize
    alt_text = alt_text.replace(' .', '.').capitalize()

    return alt_text, err

def qwen_vl(image_file_path: str, settings: dict, for_notes:bool=False, debug:bool=False) -> Tuple[str, bool]:
    """ get image description from Qwen-VL """
    err:bool = False

    # check if readonly
    image_file_path, readonly, msg = check_readonly_formats(image_file_path, settings)
    if readonly:
        return msg, False

    # prompt
    if for_notes:
        prompt:str = settings["prompt_notes"]
    else:
        prompt:str = settings["prompt"]

    model:str = settings["qwen-vl-model"]
    tokenizer:str = settings["qwen-vl-tokenizer"]

    # with Image.open(image_file_path).convert('RGB') as img:
    #     # resize image
    #     img = resize(img, settings)
    #     # prompt
    #     prompt:str = settings["prompt"]
    #     model:str = settings["cogvlm-model"]
    #     tokenizer:str = settings["cogvlm-tokenizer"]
    #     query = tokenizer.from_list_format([
    #         {'image': img},
    #         {'text': prompt},
    #     ])
    #     alt_text, history = model.chat(tokenizer, query=query, history=None)

    alt_text, history = model.chat(tokenizer, query=f'<img>{image_file_path}</img>{prompt}', history=None)

    return alt_text, err

def cog_vlm(image_file_path: str, settings: dict, for_notes:bool=False, debug:bool=False) -> Tuple[str, bool]:
    """ get image description from CogVLM """
    err:bool = False

    # check if readonly
    image_file_path, readonly, msg = check_readonly_formats(image_file_path, settings)
    if readonly:
        return msg, False

    with Image.open(image_file_path).convert('RGB') as img:

        # resize image
        img = resize(img, settings)

        # prompt
        if for_notes:
            prompt:str = settings["prompt_notes"]
        else:
            prompt:str = settings["prompt"]

        model:str = settings["cogvlm-model"]
        tokenizer:str = settings["cogvlm-tokenizer"]

        inputs = model.build_conversation_input_ids(tokenizer, query=prompt, history=[], images=[img])

    if settings["cuda_available"]:
        inputs = {
            'input_ids': inputs['input_ids'].unsqueeze(0).to('cuda'),
            'token_type_ids': inputs['token_type_ids'].unsqueeze(0).to('cuda'),
            'attention_mask': inputs['attention_mask'].unsqueeze(0).to('cuda'),
            'images': [[inputs['images'][0].to('cuda').to(torch.bfloat16)]],
        }
    # elif settings["mps_available"]:
    #     inputs = {
    #         'input_ids': inputs['input_ids'].unsqueeze(0).to('mps'),
    #         'token_type_ids': inputs['token_type_ids'].unsqueeze(0).to('mps'),
    #         'attention_mask': inputs['attention_mask'].unsqueeze(0).to('mps'),
    #         'images': [[inputs['images'][0].to('mps').to(torch.bfloat16)]],
    #     }
    else:
        inputs = {
            'input_ids': inputs['input_ids'].unsqueeze(0).to('cpu'),
            'token_type_ids': inputs['token_type_ids'].unsqueeze(0).to('cpu'),
            'attention_mask': inputs['attention_mask'].unsqueeze(0).to('cpu'),
            'images': [[inputs['images'][0].to('cpu').to(torch.bfloat16)]],
        }

    gen_kwargs = {"max_length": 2048, "do_sample": False}

    alt_text:str = ""
    with torch.no_grad():
        outputs = model.generate(**inputs, **gen_kwargs)
        outputs = outputs[:, inputs['input_ids'].shape[1]:]

        alt_text = tokenizer.decode(outputs[0])

    return alt_text, err

def img_file_to_base64(image_file_path:str , settings: dict, debug:bool=False) -> str:
    """ load image, resize, and convert to base64_str """
    with Image.open(image_file_path) as original_img:
        im = original_img.convert("RGB")

        # resize image
        im = resize(im, settings)

        # check
        buffer = io.BytesIO()
        im.save(buffer, format=original_img.format.upper())
        buffer.seek(0)

    # Encode the image bytes to Base64
    base64_bytes = base64.b64encode(buffer.getvalue())

    # str
    base64_str = base64_bytes.decode('utf-8')

    return base64_str

def gpt4v(image_file_path: str, extension:str, settings: dict, for_notes:bool=False, debug:bool=False) -> Tuple[str, bool]:
    """ get image description from GPT-4V """
    err:bool = False
    alt_text:str = "Error"

    # check if readonly
    image_file_path, readonly, msg = check_readonly_formats(image_file_path, settings)
    if readonly:
        return msg, False

    api_key = os.environ.get("OPENAI_API_KEY")
    if api_key is None or api_key == "":
        print("OPENAI_API_KEY not found in environment")
    else:
        # convert image to JPEG
        basename:str = os.path.basename(image_file_path).split(".")[0]
        jpeg_image_file_path = os.path.join(os.path.dirname(image_file_path), f"{basename}.jpg")

        with Image.open(image_file_path) as img:
            if img.format != 'JPEG':
                # Convert the image to RGB mode in case it's not
                img = img.convert('RGB')
                # Save the image as JPEG
                img.save(jpeg_image_file_path, 'JPEG')

                image_file_path = jpeg_image_file_path

        print(f"Image file size: {os.path.getsize(image_file_path)}")

        # get image and convert to base64_str
        img_base64_str = img_file_to_base64(image_file_path, settings)

        # prompt
        if for_notes:
            prompt:str = settings["prompt_notes"]
        else:
            prompt:str = settings["prompt"]

        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {api_key}"
        }
        payload = {
            "model": settings["gpt4v_model"],
            "messages": [
            {
                "role": "user",
                "content": [
                {
                    "type": "text",
                    "text": prompt
                },
                {
                    "type": "image_url",
                    "image_url": {
                    "url": f"data:image/{extension};base64,{img_base64_str}"
                    }
                }
                ]
            }
            ],
            "max_tokens": 300
        }

        gpt4v_server = "https://api.openai.com/v1/chat/completions"
        try:
            response = requests.post(gpt4v_server, headers=headers, json=payload, timeout=20)

            json_out = response.json()

            if debug:
                print(json.dumps(json_out, indent=4))

            if 'error' in json_out:
                print()
                print(json_out['error']['message'])
                err = True
            else:
                alt_text = json_out["choices"][0]["message"]["content"]
        except requests.exceptions.ConnectionError:
            print(f"ConnectionError: Unable to access the server at: '{gpt4v_server}'")
            err = True
        except TimeoutError:
            print("TimeoutError")
            err = True
        except Exception as e:
            print(f"Exception: '{str(e)}'")
            err = True

    return alt_text, err

def check_ollama_model_available(settings:dict) -> bool:
    " Check if model is available on Ollama server "
    err:bool = False
    model_specified = settings["model"]
    if ":" not in model_specified:
        model_specified = f"{model_specified}:latest"

    # check if model available
    ollama_url = f"{settings['ollama_url']}/api/tags"
    try:
        response = requests.get(ollama_url, timeout=10)
        response.raise_for_status()
        
    except requests.exceptions.ConnectionError:
        print(f"ConnectionError: Unable to access the server at: '{ollama_url}'")
        err = True
    except TimeoutError:
        print("TimeoutError")
        err = True
    else:
        json_out = response.json()
        ollama_model_response = json_out["models"]

        err = True
        all_models = []
        for m in ollama_model_response:
            if model_specified == m['name']:
                err = False
            all_models.append(m['name'])
        
        if err:
            print("Models available on the Ollama server:")
            for m in all_models:
                print(f"  {m}")
            print()
            print(f"Model '{model_specified}' not found on Ollama server: '{ollama_url}'.")
            print("Please pull the model using Ollama or use one of the other models available.\n")

    return err, model_specified

def use_ollama(image_file_path: str, extension:str, settings: dict, for_notes:bool=False, debug:bool=False) -> Tuple[str, bool]:
    """ get image description from model accessed Ollama server """
    err:bool = False
    alt_text:str = "Error"

    if not err:
        # check if readonly
        image_file_path, readonly, msg = check_readonly_formats(image_file_path, settings)
        if readonly:
            return msg, False

        # convert image to JPEG
        basename:str = os.path.basename(image_file_path).split(".")[0]
        jpeg_image_file_path = os.path.join(os.path.dirname(image_file_path), f"{basename}.jpg")

        img_base64_str = ""
        with Image.open(image_file_path) as img:
            if img.format != 'JPEG':
                # Convert the image to RGB mode in case it's not
                img = img.convert('RGB')
                # Save the image as JPEG
                img.save(jpeg_image_file_path, 'JPEG')

                image_file_path = jpeg_image_file_path

            print(f"Image file size: {os.path.getsize(image_file_path)}")

            # get image and convert to base64_str
            img_base64_str = img_file_to_base64(image_file_path, settings)

    if not err and len(img_base64_str) > 0:
        # prompt
        if for_notes:
            prompt:str = settings["prompt_notes"]
        else:
            prompt:str = settings["prompt"]

        headers = {
            "Content-Type": "application/json",
        }

        payload = {
            "model": settings["model"],
            "prompt": f"{prompt}",
            "stream": False,
            "images": [ img_base64_str ]
        }

        # http://localhost:11434/api/generate
        ollama_url = f"{settings['ollama_url']}/api/generate"
        try:
            response = requests.post(ollama_url, headers=headers, json=payload, timeout=60)
            response.raise_for_status()
        except requests.exceptions.ConnectionError:
            print(f"ConnectionError: Unable to access the server at: '{ollama_url}'")
            err = True
        except TimeoutError:
            print("TimeoutError")
            err = True
        else:
            json_out = response.json()

            if debug:
                print(json.dumps(json_out, indent=4))

            if 'error' in json_out:
                print("ERROR in ouput")
                print(json.dumps(json_out, indent=4))
                err = True
            else:
                alt_text = json_out["response"]
                # remove newlines
                alt_text = alt_text.replace("\n", " ")
                # remove double spaces
                alt_text = alt_text.replace("  ", " ")
    else:
        print("Error, image size is zero")
        err = True

    return alt_text, err

def accessibility_report(out_file_name: str, pptx_file_name: str, debug:bool = False) -> None:
    """
    Create accessibility report based on infomation in the text file generated
    """
    # accessibility report
    print("---- Accessibility report --------------------------------------------")
    print(f"Powerpoint file: '{pptx_file_name}'")
    empty_alt_txt:int = 0
    alt_text_list:list = []
    img_cnt:int = 0
    with open(out_file_name, "r", encoding="utf-8") as file:
        # tab delimited file
        csv_reader = csv.reader(file, delimiter="\t")

        # skip header
        next(csv_reader)

        # process rows
        for row in csv_reader:
            if len(row) == 10 and len(row[8]) > 0 and not str2bool(row[7]):
                # not decorative
                if len(row[6]) == 0:
                    if debug: print(row)
                    empty_alt_txt += 1
                if row[4] == "Picture":
                    img_cnt += 1

                # create list of alt text length
                alt_text_list.append(int(row[7]))
            elif len(row) != 10:
                print(f"Unexpected row length: {len(row)}, row: {row}")

    print(f"Images: {img_cnt}")
    print(f"Objects: {csv_reader.line_num - 1}")

    print(f"Number of missing alt texts for Group(s), Image(s) or Objects(s): {empty_alt_txt}")
    print(f"Min alt text length: {min(alt_text_list)}")
    print(f"Max alt text length: {max(alt_text_list)}")

    print("----------------------------------------------------------------------")


def replace_alt_texts(file_path: str, file_path_txt_file: str, debug:bool = False) -> bool:
    """
    Replace alt texts specified in a text file (e.g. generated by this script and edited to correct or improve)
    Text file should have a header and the same columns as the output files generated by this script
    """
    err:bool = False

    # Check if text file is exists
    if not os.path.isfile(file_path_txt_file):
        print(f"Unable to access file: {file_path_txt_file}")
        return False

    # get name, extension, folder from Powerpoint file
    name:str = pathlib.Path(file_path).stem
    extension:str = pathlib.Path(file_path).suffix
    dirname:str = os.path.dirname(file_path)

    # process txt file
    print(f"Reading: {file_path_txt_file}...")
    csv_rows:list[str] = []
    with open(file_path_txt_file, "r", encoding="utf-8") as file:
        # assume tab delimited file
        csv_reader = csv.reader(file, delimiter="\t")

        # skip header
        next(csv_reader)

        for row in csv_reader:
            csv_rows.append(row)

    # process powerpoint file
    print(f"Processing Powerpoint file: {file_path}")
    prs = Presentation(file_path)

    # Loop through slides
    object_cnt:int = 0
    for slide_cnt, slide in enumerate(prs.slides):
        # loop through shapes
        slide_object_cnt = 0
        for shape in slide.shapes:
            _, object_cnt, slide_object_cnt = process_shapes_from_file(shape, None, csv_rows, slide_cnt, slide_object_cnt, object_cnt, debug)

    if not err:
        # Save file
        outfile:str = os.path.join(dirname, f"{name}_alt_text{extension}")
        print(f"Saving Powerpoint file with new alt-text to: '{outfile}'")
        prs.save(outfile)

    return err

def process_shapes_from_file(shape: BaseShape, group_shape_list: list[BaseShape], csv_rows, slide_cnt:int, slide_object_cnt:int, object_cnt: int, debug:bool) -> int:
    """ recursive function to process shapes and shapes within groups """
    # Check if the shape has a picture
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        if group_shape_list is None:
            group_shape_list = [shape]
        else:
            group_shape_list.append(shape)

        for embedded_shape in shape.shapes:
            group_shape_list, object_cnt, slide_object_cnt = process_shapes_from_file(embedded_shape, group_shape_list, csv_rows, slide_cnt, slide_object_cnt, object_cnt, debug)

        # current group shape (last one)
        group_shape = group_shape_list[-1]

        # get decorative
        decorative_pptx:bool = is_decorative(group_shape)
        decorative:bool = str2bool(csv_rows[object_cnt][8])

        # change decorative status
        if decorative_pptx != decorative:
            # set decorative status of image
            print(f"Side: {slide_cnt}, {group_shape.name}, can't set the docorative status to: {bool2str(decorative)}")

        alt_text: str = ""
        if not decorative:
            # get alt text from text file
            # print(f"Set to {csv_rows[image_cnt][6]}")
            alt_text = csv_rows[object_cnt][6]

        # set alt text
        if debug: print(f"Set group to {alt_text}")
        set_alt_text(group_shape, alt_text)

        slide_object_cnt += 1
        object_cnt += 1

        # remove last one
        group_shape_list = group_shape_list[:-1]

    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:

        # get decorative
        decorative_pptx:bool = is_decorative(shape)
        decorative:bool = str2bool(csv_rows[object_cnt][8])

        # change decorative status
        if decorative_pptx != decorative:
            # set decorative status of image
            print(f"Side: {slide_cnt}, {shape.name}, can't set the docorative status to: {bool2str(decorative)}")

        alt_text: str = ""
        if not decorative:
            # get alt text from text file
            alt_text = csv_rows[object_cnt][6]

        # set alt text
        set_alt_text(shape, alt_text)

        slide_object_cnt += 1
        object_cnt += 1

    elif shape.shape_type in [MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.FREEFORM, \
                              MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.IGX_GRAPHIC, MSO_SHAPE_TYPE.CANVAS, \
                              MSO_SHAPE_TYPE.MEDIA, MSO_SHAPE_TYPE.WEB_VIDEO, MSO_SHAPE_TYPE.DIAGRAM, \
                              MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT, MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT, \
                              MSO_SHAPE_TYPE.WEB_VIDEO, MSO_SHAPE_TYPE.LINKED_OLE_OBJECT, \
                              MSO_SHAPE_TYPE.CALLOUT]:

        # get decorative
        decorative_pptx:bool = is_decorative(shape)
        decorative:bool = str2bool(csv_rows[object_cnt][8])

        # change decorative status
        if decorative_pptx != decorative:
            # set decorative status of image
            print(f"Side: {slide_cnt}, {shape.name}, can't set the docorative status to: {bool2str(decorative)}")

        alt_text: str = ""
        if not decorative:
            # get alt text from text file
            alt_text = csv_rows[object_cnt][6]

        # set alt text
        set_alt_text(shape, alt_text)

        slide_object_cnt += 1
        object_cnt += 1

    elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:

        decorative_pptx:bool = is_decorative(shape)
        decorative:bool = str2bool(csv_rows[object_cnt][8])

        # change decorative status
        if decorative_pptx != decorative:
            # set decorative status of image
            print(f"Side: {slide_cnt}, {shape.name}, can't set the docorative status to: {bool2str(decorative)}")

        alt_text: str = ""
        if not decorative:
            # get alt text from text file
            alt_text = csv_rows[object_cnt][6]

        # set alt text
        set_alt_text(shape, alt_text)

    return group_shape_list, object_cnt, slide_object_cnt

def add_presenter_note(pptx_path:str, pptx:dict, settings:dict) -> bool:
    """ add presenter note to each slide """
    err:bool = False

    pptx_name:str = pathlib.Path(pptx_path).stem
    #pptx_extension:str = pathlib.Path(pptx_path).suffix
    dirname:str = os.path.dirname(pptx_path)

    slide_dir:str = os.path.join(dirname, pptx_name, "slides_png")
    current_slide:int = pptx["slide_cnt"]
    t = num2str(pptx['pptx_nslides'], current_slide + 1)
    file_name:str = f"{pptx_name}-{t}.png"
    slide_image_file_path = os.path.join(slide_dir, file_name)

    if os.path.isfile(slide_image_file_path):
        alt_text, err = generate_description(slide_image_file_path, ".png", settings, for_notes=True)
        slide = pptx["current_slide"]
        slide.notes_slide.notes_text_frame.text = alt_text
        print(f"Slide: {current_slide + 1}\t{alt_text}")
    else:
        print(f"Unable to access image file: {slide_image_file_path}")
        err = True

    return err

def remove_presenter_notes(pptx_path:str) -> bool:
    """ remove all presenter notes """
    err:bool = False

    # get name, extension, folder from Powerpoint file
    pptx_name:str = pathlib.Path(pptx_path).stem
    pptx_extension:str = pathlib.Path(pptx_path).suffix
    dirname:str = os.path.dirname(pptx_path)

    # process powerpoint file
    print(f"Processing Powerpoint file: {pptx_path}")
    prs = Presentation(pptx_path)

    # Loop through slides
    for _, slide in enumerate(prs.slides):
        slide.notes_slide.notes_text_frame.text = ""

    new_pptx_file_name = os.path.join(dirname, f"{pptx_name}_notes_removed{pptx_extension}")
    prs.save(new_pptx_file_name)
    print(f"Saved Powerpoint file with presenter notes removed to: '{new_pptx_file_name}'\n")

    return err

def export_presenter_notes(pptx_path:str) -> bool:
    """ export presenter notes """
    err:bool = False

    # get name and folder from Powerpoint file
    pptx_name:str = pathlib.Path(pptx_path).stem
    dirname:str = os.path.dirname(pptx_path)

    # output file
    notes_file_path = os.path.join(dirname, f"{pptx_name}_notes.txt")

    with open(notes_file_path, "w", encoding="utf-8") as out_file:
        print(f"Processing Powerpoint file: {pptx_path}")
        prs = Presentation(pptx_path)

        # Loop through slides
        for slide_cnt, slide in enumerate(prs.slides):
            if slide.shapes.title:
                title = slide.shapes.title.text
            else:
                title = ""
            str = f"=== Slide {slide_cnt} - {title} ===\n\n{slide.notes_slide.notes_text_frame.text}\n\n"
            out_file.write(str)

    if not err:
        print(f"Exported presenter notes to file: '{notes_file_path}'")

    return err

def export_slides_to_images(pptx_path:str) -> bool:
    """ export slides to PNG, Windows ONLY and requires that Powerpoint is installed """

    err:bool = False
    dirname:str = os.path.dirname(pptx_path)
    pptx_name:str = pathlib.Path(pptx_path).stem

    prs:Presentation = Presentation(pptx_path)
    n_slides = len(prs.slides)

    # create folder to store images
    img_folder = os.path.join(dirname, pptx_name)
    if not os.path.isdir(img_folder):
        os.makedirs(img_folder)

    path_to_folder_to_save = os.path.join(dirname, pptx_name, "slides_png")
    if not os.path.isdir(path_to_folder_to_save):
        os.makedirs(path_to_folder_to_save)

    if platform.system() == "Windows":
        try:
            import comtypes.client

            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            powerpoint.Visible = 1

            abs_file_path = os.path.abspath(pptx_path)
            presentation = powerpoint.Presentations.Open(abs_file_path)

            abs_path_to_folder_to_save = os.path.abspath(path_to_folder_to_save)

            for i, slide in enumerate(presentation.Slides):
                p = num2str(n_slides, i + 1)
                file_path = os.path.join(abs_path_to_folder_to_save, f"{pptx_name}-{p}.png")
                slide.Export(file_path, "PNG")

            presentation.Close()
            powerpoint.Quit()
            print(f"Slides saved as PNG images in folder: '{abs_path_to_folder_to_save}'")
        except Exception as e:
            print(f"Unable to export slides: {str(e)}")
            err = True

    elif platform.system() != "Windows":
        # export PPTX first to PDF
        print("Exporting to PDF...")

        cmd:list[str] = ["soffice", "--headless", "--convert-to", "pdf", pptx_path, "--outdir", path_to_folder_to_save]
        path_to_cmd = shutil.which(cmd[0])
        if path_to_cmd is not None:
            r = subprocess.run(cmd, stdin=subprocess.PIPE, stdout=subprocess.PIPE, shell=False, check=True)
        else:
            print("Warning, LibreOffice not installed.")

        # save each page as a separate file
        print("Extracting pages from PDF...")
        pdf_file:str = os.path.join(path_to_folder_to_save, f"{pptx_name}.pdf")
        out_file_name:str = os.path.join(path_to_folder_to_save, f"{pptx_name}.pdf")

        cmd = ["qpdf", "--split-pages", pdf_file, out_file_name]
        path_to_cmd = shutil.which(cmd[0])
        if path_to_cmd is not None:
            r = subprocess.run(cmd, stdin=subprocess.PIPE, stdout=subprocess.PIPE, shell=False, check=True)
        else:
            print("Warning, qpdf not installed.")

        # export from PDF to PNG
        print("Converting each page to PNG...")
        if platform.system() == "Darwin":
            the_files = os.listdir(path_to_folder_to_save)
            for f in the_files:
                if not f.startswith(".") and f.startswith(f"{pptx_name}-"):
                    in_file = os.path.join(path_to_folder_to_save, f)
                    cmd = ["sips", "-s", "dpiWidth", "300", "-s", "dpiHeight", "300", "-s", "format", "png", in_file, "--out", path_to_folder_to_save]
                    path_to_cmd = shutil.which(cmd[0])

                    if path_to_cmd is not None:
                        r = subprocess.run(cmd, stdin=subprocess.PIPE, stdout=subprocess.PIPE, shell=False, check=True) #capture_output = True)
                    else:
                        print("Unable to find 'sips'")
        else:
            # Linux
            the_files = os.listdir(path_to_folder_to_save)
            for f in the_files:
                if not f.startswith(".") and f.startswith(f"{pptx_name}-"):
                    in_file = os.path.join(path_to_folder_to_save, f"{f}[0]")
                    out_file = f"{pathlib.Path(in_file).stem}.png"
                    cmd = ["convert", in_file, "-density", "300", out_file]
                    path_to_cmd = shutil.which(cmd[0])

                    if path_to_cmd is not None:
                        r = subprocess.run(cmd, stdin=subprocess.PIPE, stdout=subprocess.PIPE, shell=False, check=True) #capture_output = True)
                    else:
                        print("Unable to find 'convert'")
        # remove PDFs
        for f in the_files:
            if not f.startswith(".") and f.startswith(f"{pptx_name}-"):
                in_file = os.path.join(path_to_folder_to_save, f)
                pid = os.remove(in_file)
                if pid:
                    _ = os.wait()

        # remove exported PDF file
        #os.remove(pdf_file)

        print(f"Slides saved as PNG images in folder: '{path_to_folder_to_save}'")
    else:
        err = True
        print("Unable to convert PPTX to images.")

    #if not err:
    #    # generate pptx from slide images
    #    images_to_pptx(path_to_folder_to_save, f"{pptx_name}_slides.pptx")

    return err

def images_to_pptx(images_path:str, output_file:str="presentation.pptx", debug:bool=False) -> None:
    """
    Create a PowerPoint presentation with an image on each slide
    """
    print("Generate pptx from images...")

    # Create a presentation object
    prs = Presentation()

    # Define slide width and height (in centimeters)
    slide_width_cm = prs.slide_width.cm
    slide_height_cm = prs.slide_height.cm

    # select only PNG files
    the_files = os.listdir(images_path)
    the_files = [f for f in the_files if f.endswith(".png")]
    if len(the_files) == 0:
        print("No images found!")
        return

    for f in the_files:
        file_path = os.path.join(images_path, f)

        # Open the image to get its size
        with Image.open(file_path) as img:
            width_px, height_px = img.size
            dpi = img.info.get('dpi', (96, 96))  # Defaulting to 96 DPI if not provided

        # Calculate the image size in centimeters
        width_cm = width_px / dpi[0] * 2.54
        height_cm = height_px / dpi[1] * 2.54

        # Calculate the scaling factor
        scale_factor = min(slide_width_cm / width_cm, slide_height_cm / height_cm)

        # Add a slide with a blank layout
        slide_layout = prs.slide_layouts[5]  # 5 is the index for a blank slide
        slide = prs.slides.add_slide(slide_layout)

        # Add and resize the image to the slide
        img = slide.shapes.add_picture(file_path, Cm(0), Cm(0),
                                       width=Cm(width_cm * scale_factor),
                                       height=Cm(height_cm * scale_factor))

    # Save the presentation
    out = os.path.join(images_path, output_file)
    prs.save(out)

    print(f"Saved pptx to '{out}'.")

# argv: List[str]
def main() -> int:
    """ main """
    err:bool = False

    parser = argparse.ArgumentParser(description='Add alt-text automatically to images and objects in Powerpoint')
    parser.add_argument("file", type=str, help="Powerpoint file")
    parser.add_argument("--report", action='store_true', default=False, help="flag to generate alt text report")
    parser.add_argument("--model", type=str, default="", help="kosmos-2, openclip, llava, gpt-4v")

    # Ollama
    parser.add_argument("--use_ollama", action='store_true', default=False, help="use Ollama server")
    parser.add_argument("--server", type=str, default="http://localhost", help="Ollama server URL, default=http://localhost")
    parser.add_argument("--port", type=str, default="11434", help="Ollama server port, default=11434")

    # OpenCLIP
    parser.add_argument("--show_openclip_models", action='store_true', default=False, help="show OpenCLIP models and pretrained models")
    parser.add_argument("--openclip_model", type=str, default="coca_ViT-L-14", help="OpenCLIP model")
    parser.add_argument("--openclip_pretrained", type=str, default="mscoco_finetuned_laion2B-s13B-b90k", help="OpenCLIP pretrained model")
    #
    parser.add_argument("--resize", type=str, default="500", help="resize image to same width and height in pixels, default:500, use 0 to disable resize")
    #
    parser.add_argument("--prompt", type=str, default="", help="custom prompt")
    parser.add_argument("--prompt_notes", type=str, default="", help="custom prompt for presenter notes")
    #
    #parser.add_argument("--save", action='store_true', default=False, help="flag to save powerpoint file with updated alt texts")
    parser.add_argument("--replace", type=str, default="", help="replace alt texts in pptx with those specified in file")
    parser.add_argument("--remove_presenter_notes", action='store_true', default=False, help="remove all presenter notes")
    parser.add_argument("--export_presenter_notes", action='store_true', default=False, help="export presenter notes")
    parser.add_argument("--export_slides", action='store_true', default=False, help="export pptx slides to png images")
    #
    parser.add_argument("--add_to_notes", action='store_true', default=False, help="add image descriptions to slide notes")
    #
    parser.add_argument("--debug", action='store_true', default=False, help="flag for debugging")

    args = parser.parse_args()

    prompt:str = args.prompt
    model_str:str = args.model.lower()

    if args.show_openclip_models:
        openclip_models = open_clip.list_pretrained()
        print("OpenCLIP models:")
        for m, p in openclip_models:
            print(f"Model: {m}, pretrained model: {p}")
        return int(err)

    # set default prompt
    if model_str == "gpt-4v":
        if args.prompt == "":
            prompt = "Describe the image using one or two sentences. Do not mention the word 'image'."
    elif model_str == "kosmos-2":
        if args.prompt == "":
            prompt = "<grounding>An image of"
            #prompt = "<grounding>Describe this image:"
    elif model_str == "qwen-vl":
        if args.prompt == "":
            prompt = "Describe the image using one or two sentences."
    elif model_str == "cogvlm":
        if args.prompt == "":
            prompt = "Describe the image using one or two sentences."
    elif args.use_ollama:
        if args.prompt == "":
            prompt = "You are an expert at understanding images and graphs. Answer concisely for someone who is visually impaired. Describe what you see. Your response should be one or two sentences."
    else:
        if args.prompt == "":
            prompt = "Describe image. Your response should be one or two sentences."

    if args.prompt_notes == "":
        prompt_presenter_notes = "Describe the image in a few sentences for someone who is visually impaired. Start the desciption with 'This slide'"
    else:
        prompt_presenter_notes = args.prompt_notes

    # Read PowerPoint file and list images
    powerpoint_file_name = args.file
    if not os.path.isfile(powerpoint_file_name):
        print(f"Error: File {powerpoint_file_name} not found.")
        err = True
    else:
        settings:dict = {
            "report": args.report,
            "model": model_str,
            "kosmos2_model": None,
            "kosmos2_pretrained": None,
            "openclip_model_name": args.openclip_model,
            "openclip_pretrained": args.openclip_pretrained,
            "openclip-model": None,
            "openclip-preprocess": None,
            "qwen-vl-model": None,
            "qwen-vl-tokenizer": None,
            "cogvlm-model": None,
            "cogvlm-tokenizer": None,
            "use_ollama": args.use_ollama,
            "ollama_url": f"{args.server}:{args.port}",
            "gpt4v_model": "gpt-4-vision-preview",
            "cuda_available": torch.cuda.is_available(),
            "mps_available": torch.backends.mps.is_available(),
            "prompt": prompt,
            "prompt_notes": prompt_presenter_notes,
            "img_size": int(args.resize),
            "add_to_notes": args.add_to_notes
        }

        if args.replace != "":
            # file with alt text provided
            err = replace_alt_texts(powerpoint_file_name, args.replace, args.debug)
        elif args.remove_presenter_notes:
            err = remove_presenter_notes(powerpoint_file_name)
        elif args.export_presenter_notes:
            err = export_presenter_notes(powerpoint_file_name)
        elif args.export_slides:
            err = export_slides_to_images(powerpoint_file_name)
        else:
            if args.add_to_notes:
                print(f"Model: {model_str}")
                print(f"Presenter notes prompt: '{prompt_presenter_notes}'")
                # export slides to images so that model can interpret the whole slide
                err = export_slides_to_images(powerpoint_file_name)

            # add alt-text
            if not err:
                err = process_images_from_pptx(powerpoint_file_name, settings, args.debug)

    return int(err)

if __name__ == "__main__":
    EXIT_CODE = main() #main(sys.argv[1:])
    sys.exit(EXIT_CODE)
